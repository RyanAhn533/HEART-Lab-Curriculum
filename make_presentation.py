# -*- coding: utf-8 -*-
"""
딥러닝실습 발표 PPT - 템플릿 기반 재설계
템플릿의 마스터 레이아웃/배경/폰트를 그대로 활용
"""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

IMG = 'C:/Users/Ryan/HEART-Lab-Curriculum/images/'

# 템플릿 기반 시작 — 기존 슬라이드 전부 삭제 후 레이아웃만 사용
prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/template.pptx')

# 기존 슬라이드 전부 삭제
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldId = prs.slides._sldIdLst[0]
    prs.slides._sldIdLst.remove(sldId)

# 색상
NAVY = RGBColor(0x01, 0x3B, 0x94)
SKY = RGBColor(0x86, 0xB7, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DGRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF0, 0xF2, 0xF5)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x00, 0xA8, 0x6B)
RED = RGBColor(0xE0, 0x3E, 0x3E)
FONT = 'Noto Sans KR'

# Layout references
LY_TITLE = prs.slide_layouts[1]    # 제목 슬라이드
LY_CONTENT = prs.slide_layouts[0]  # 1_제목 및 내용
LY_BLANK = prs.slide_layouts[7]    # 빈 화면
LY_SECTION = prs.slide_layouts[3]  # 구역 머리글
LY_TITLEONLY = prs.slide_layouts[6]  # 제목만

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FONT; p.alignment = align
    return tb

def lines(slide, l, t, w, h, items, sz=16, color=BLACK, sp=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        b = False
        if item.startswith('*'):
            item = item[1:]; b = True
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.bold = b; p.font.name = FONT; p.space_after = Pt(sp)
    return tb

def img(slide, path, l, t, width=None, height=None):
    fp = IMG + path if not os.path.isabs(path) else path
    if not os.path.exists(fp): return None
    kw = {}
    if width: kw['width'] = width
    if height: kw['height'] = height
    return slide.shapes.add_picture(fp, l, t, **kw)

def slide_header(slide, title, subtitle=None):
    """공통 헤더: 진한파랑 상단 바 + 제목"""
    rect(slide, 0, 0, Inches(13.333), Inches(1.1), SKY)
    rect(slide, 0, Inches(1.1), Inches(13.333), Inches(0.04), NAVY)
    txt(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
        title, 28, NAVY, True)
    if subtitle:
        txt(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.3),
            subtitle, 13, GRAY)
    # 하단 바
    rect(slide, 0, Inches(7.2), Inches(13.333), Inches(0.3), NAVY)
    txt(slide, Inches(0.4), Inches(7.2), Inches(5), Inches(0.3),
        'HEART Lab | Sejong University', 9, WHITE)

# ============================================================
# SLIDE 1: 타이틀
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
rect(s, 0, 0, Inches(13.333), Inches(7.5), SKY)
rect(s, 0, Inches(2.2), Inches(13.333), Inches(3.2), NAVY)
rect(s, 0, Inches(7.1), Inches(13.333), Inches(0.4), NAVY)
txt(s, Inches(0.8), Inches(0.6), Inches(5), Inches(0.5),
    '2026-1학기 딥러닝실습', 16, NAVY)
txt(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.0),
    '딥러닝 실습, 왜 하는 걸까?', 44, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
    '대학, 취업, 그리고 AI', 20, RGBColor(0xBB,0xD5,0xFF), False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
    '안준영  |  세종대학교 HEART Lab', 18, NAVY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: 고등학교 vs 대학교
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '고등학교 vs 대학교')

# 고등학교 카드
rrect(s, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.8), LGRAY)
txt(s, Inches(0.6), Inches(1.6), Inches(5.5), Inches(0.6),
    '고등학교', 26, GRAY, True, PP_ALIGN.CENTER)
lines(s, Inches(1.2), Inches(2.4), Inches(4.3), Inches(2.5), [
    '교복을 입는다',
    '선생님이 다 챙겨준다',
    '보호받는 공간',
    '',
    '*→ 미성년자',
], 18, GRAY, 10)

# 화살표
txt(s, Inches(6.0), Inches(3.0), Inches(1.3), Inches(0.8),
    '→', 44, NAVY, True, PP_ALIGN.CENTER)

# 대학교 카드
rrect(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.8), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.6),
    '대학교', 26, NAVY, True, PP_ALIGN.CENTER)
lines(s, Inches(7.8), Inches(2.4), Inches(4.3), Inches(2.5), [
    '교복이 없다',
    '스스로 해야 한다',
    '보호가 없다',
    '',
    '*→ 성인',
], 18, NAVY, 10)

# 결론
rrect(s, Inches(2.5), Inches(5.8), Inches(8.3), Inches(0.8), NAVY)
txt(s, Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.6),
    '대학교는 결국 취업을 준비하는 공간', 24, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: 이 수업 왜?
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '그래서 이 수업을 왜 듣는가?')

txt(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.6),
    '딥러닝 실습 — 코드를 돌리고, 프로젝트를 하고', 24, DGRAY)
txt(s, Inches(0.8), Inches(2.8), Inches(11), Inches(0.6),
    '학점?  재미?', 28, GRAY)

rrect(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.4), NAVY)
txt(s, Inches(1.2), Inches(4.1), Inches(11), Inches(0.6),
    '"나는 이걸 할 수 있습니다"를 증명하는 과정', 28, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(4.8), Inches(11), Inches(0.5),
    '이 수업의 프로젝트 = 취업 준비', 20, RGBColor(0xBB,0xD5,0xFF), False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: AI 취업 두 가지
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '인공지능 취업, 크게 두 가지')

# 기업 개발자
rrect(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.5), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.5),
    '기업 개발자', 24, NAVY, True, PP_ALIGN.CENTER)
lines(s, Inches(1.0), Inches(2.2), Inches(5.0), Inches(2.0), [
    '각 회사의 도메인이 다름',
    '(반도체, 자동차, 의학, 금융, 로봇...)',
    '',
    '도메인에 맞는 서비스(제품)를 만드는 사람',
], 15, DGRAY, 6)

# 기업 로고
img(s, 'samsung_logo.png', Inches(1.0), Inches(4.3), height=Inches(0.5))
img(s, 'kakao_logo.jpg', Inches(3.0), Inches(4.2), height=Inches(0.6))

# 연구소
rrect(s, Inches(6.9), Inches(1.4), Inches(5.8), Inches(3.5), RGBColor(0xFE,0xF0,0xE8))
txt(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.5),
    '연구소 (연구원)', 24, ORANGE, True, PP_ALIGN.CENTER)
lines(s, Inches(7.3), Inches(2.2), Inches(5.0), Inches(2.0), [
    '기업 연구소, 국책 연구소, 대학원',
    '트렌드에 맞춰 기술 발전을 리드',
    '논문, 특허, 기술 이전',
    '',
    '*연구소는 국책과제(정부 과제)로 운영',
], 15, DGRAY, 6)

# NVIDIA 로고
img(s, 'nvidia_logo.png', Inches(9.0), Inches(4.2), height=Inches(0.6))

# 하단
rrect(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.7), NAVY)
txt(s, Inches(1.5), Inches(5.35), Inches(10.3), Inches(0.6),
    '공통:  AI로 문제를 풀 수 있는 사람이 필요하다', 20, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: 필요한 능력
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '우리에게 중요한 능력은?')

rrect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.0), NAVY)
txt(s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.8),
    '서비스 = Question (문제)  →  "이 문제를 어떻게 풀 수 있는가?"', 24, WHITE, True, PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.4),
    '어떤 데이터로,  어떤 모델로,  어떤 기술로?', 20, DGRAY)

# 기술 카드 4개
techs = [
    ('LLM', 'ChatGPT, Claude\n대화형 AI', NAVY),
    ('AI Agent', '자율적으로 일하는 AI\n도구 사용, 의사결정', RGBColor(0x3B,0x82,0xF6)),
    ('멀티모달', '영상 + 음성 + 텍스트\n여러 데이터를 결합', ORANGE),
    ('Physical AI', '로봇, 자율주행\n물리 세계의 AI', GREEN),
]
for i, (title, desc, c) in enumerate(techs):
    left = Inches(0.6 + i * 3.15)
    rrect(s, left, Inches(3.4), Inches(2.9), Inches(2.5), c)
    txt(s, left, Inches(3.6), Inches(2.9), Inches(0.5),
        title, 22, WHITE, True, PP_ALIGN.CENTER)
    txt(s, left, Inches(4.3), Inches(2.9), Inches(1.2),
        desc, 14, WHITE, False, PP_ALIGN.CENTER)

txt(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.4),
    '이걸 조합해서 실제 서비스를 구현할 수 있는 능력 = 핵심 경쟁력', 18, NAVY, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: 취업시장 현실
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '그런데 요즘 취업시장 현실')

# 통계 카드 3개
stats = [
    ('53.5%→37.4%', 'SW 신입 채용 비중\n2년 만에 16.1%p 급감'),
    ('12.4%', '2026년 전체 채용 중\n신입 비율'),
    ('21만 1천 개', 'ChatGPT 이후 3년간\n청년 일자리 증발'),
]
for i, (num, desc) in enumerate(stats):
    left = Inches(0.6 + i * 4.2)
    rrect(s, left, Inches(1.4), Inches(3.9), Inches(2.4), RED)
    txt(s, left, Inches(1.6), Inches(3.9), Inches(0.7),
        num, 30, WHITE, True, PP_ALIGN.CENTER)
    txt(s, left, Inches(2.5), Inches(3.9), Inches(1.0),
        desc, 14, WHITE, False, PP_ALIGN.CENTER)

txt(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.3),
    '출처: 서울신문 2026.01 / ZDNet 2025.12 / 한국은행', 11, GRAY)

rrect(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.6), NAVY)
lines(s, Inches(2.0), Inches(4.7), Inches(9.3), Inches(1.4), [
    '*기업이 원하는 건:  바로 투입 가능한 사람',
    '',
    '스펙이 아니라,  실제로 문제를 풀어본 경험이 있는 사람',
], 20, WHITE, 6)

# ============================================================
# SLIDE 7: Top-down 사고
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '바로 투입 가능한 사람 = Top-down 사고')

# 일론 머스크 영역
rrect(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(4.0), LGRAY)
# 사진
pic = img(s, 'elon_musk.jpeg', Inches(0.8), Inches(1.6), width=Inches(3.0))
txt(s, Inches(3.9), Inches(1.6), Inches(2.5), Inches(0.4),
    'First Principles', 20, DGRAY, True)
txt(s, Inches(3.9), Inches(2.0), Inches(2.5), Inches(0.4),
    'Thinking', 20, DGRAY, True)
txt(s, Inches(3.9), Inches(2.5), Inches(2.5), Inches(0.3),
    '— Elon Musk', 13, GRAY)
lines(s, Inches(0.9), Inches(3.4), Inches(5.5), Inches(1.8), [
    '문제를 분자 단위로 분해한다',
    '기존 방식을 의심, 근본부터 다시 생각',
    '로켓 $65M → 원자재 = 2% → 10배 절감',
], 14, DGRAY, 5)

# AI 적용
rrect(s, Inches(7.0), Inches(1.4), Inches(5.7), Inches(4.0), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(7.3), Inches(1.5), Inches(5.1), Inches(0.5),
    'AI에서도 똑같다', 22, NAVY, True)

# 계단식 흐름
steps_ai = ['문제가 뭐지?', '→ 데이터는?', '  → 전처리는?', '    → 모델은?', '      → 평가는?']
for j, step in enumerate(steps_ai):
    y = Inches(2.3 + j * 0.45)
    if j == 0:
        rrect(s, Inches(7.5), y, Inches(4.8), Inches(0.38), NAVY)
        txt(s, Inches(7.6), y, Inches(4.6), Inches(0.35), step, 15, WHITE, True)
    else:
        rrect(s, Inches(7.5), y, Inches(4.8), Inches(0.38), RGBColor(0xC0,0xD4,0xF4))
        txt(s, Inches(7.6), y, Inches(4.6), Inches(0.35), step, 14, NAVY)

txt(s, Inches(7.3), Inches(4.7), Inches(5.1), Inches(0.5),
    '이 사고가 바로 나오는 사람\n= 바로 투입 가능한 사람', 16, NAVY, True)

# 하단
rrect(s, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.7), NAVY)
txt(s, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.6),
    '연구를 위한 연구가 아니라, 문제를 정의하고 풀어내는 것', 19, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: K-MER + Jetson
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '나는 취업하려고 이걸 만들었다 — K-MER')

# K-MER 소개
rrect(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(3.2), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(0.9), Inches(1.5), Inches(7.0), Inches(0.5),
    'K-MER  |  차량 내 운전자 감정인식 시스템', 20, NAVY, True)
lines(s, Inches(0.9), Inches(2.2), Inches(6.8), Inches(2.2), [
    '*산업부 수십억 규모 국책과제',
    '',
    '카메라 + 음성 + 생체신호  →  운전자 감정 실시간 판단',
    '',
    '연구를 위한 연구 X  →  실제 차량에 실증하는 기술',
    '졸업할 때 "이거 만들었습니다" 포트폴리오',
], 15, DGRAY, 4)

# Jetson 카드
rrect(s, Inches(8.4), Inches(1.4), Inches(4.3), Inches(3.2), NAVY)
# NVIDIA 로고
img(s, 'nvidia_logo.png', Inches(9.1), Inches(1.5), width=Inches(2.4))
txt(s, Inches(8.6), Inches(2.2), Inches(3.9), Inches(0.4),
    'On-Device AI  |  Jetson', 16, SKY, True, PP_ALIGN.CENTER)
lines(s, Inches(8.7), Inches(2.8), Inches(3.8), Inches(1.6), [
    'DL 모델을 엣지 디바이스에',
    '최적화하여 배포',
    '',
    'Cloud 없이 실시간 추론',
    '(저지연 + 개인정보 보호)',
    '',
    '멀티모달 경량화 → 임베디드 동작',
], 12, WHITE, 3)

# K-MER 영상 placeholder
rrect(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.8), RGBColor(0xF8,0xF0,0xE8))
rect(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.05), ORANGE)
txt(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.8),
    '[ K-MER 시연 영상 삽입 ]', 24, ORANGE, True, PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.4),
    '→ 여기에 영상 또는 스크린샷을 넣으세요', 13, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: 트렌드 + HEART Lab
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '이걸 어디서 만들었냐')

# HEART Lab 소개
rrect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.6), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.6),
    'HEART Lab', 30, NAVY, True, PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.3),
    'Human Emotion  and  Intelligent Agent  Research  for  Future Transformation', 14, GRAY, False, PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.3),
    '사람의 감정  +  지능형 에이전트  +  미래 기술', 14, NAVY, True, PP_ALIGN.CENTER)

# "이 분야가 취업이 되는 이유"
txt(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.4),
    '이 분야가 취업이 되는 이유 — 우리 말을 NVIDIA가 반증한다', 20, DGRAY, True)

# 3개 증거 카드
proofs = [
    ('NVIDIA', '감정 인식 연구를 직접 수행\nGTC에서 Emotion AI 다룸', NAVY),
    ('Affectiva (Smart Eye)', '감정 AI 글로벌 1위\n차량 감정인식 이미 상용화\n$73.5M 인수', RGBColor(0x3B,0x82,0xF6)),
    ('Top Conference', '감정 인식 논문이\nCVPR, AAAI 등에서\n가장 활발한 주제', GREEN),
]
for i, (title, desc, c) in enumerate(proofs):
    left = Inches(0.6 + i * 4.2)
    rrect(s, left, Inches(3.9), Inches(3.9), Inches(2.4), c)
    txt(s, left, Inches(4.0), Inches(3.9), Inches(0.5),
        title, 20, WHITE, True, PP_ALIGN.CENTER)
    txt(s, left, Inches(4.7), Inches(3.9), Inches(1.4),
        desc, 14, WHITE, False, PP_ALIGN.CENTER)

# 로고 배치
img(s, 'nvidia_logo.png', Inches(1.3), Inches(5.7), height=Inches(0.5))
img(s, 'affectiva_logo.png', Inches(5.5), Inches(5.5), height=Inches(0.7))

txt(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
    '이 트렌드 위에 있는 사람이 취업이 되는 것', 18, NAVY, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: 커리큘럼
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, '이 능력을 어떻게 키우냐?')

lines(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.0), [
    '카카오 AI 부트캠프, 비트캠프 경험 → 취업 목적 교육은 구조가 다르다',
    '그 경험을 바탕으로 교수님과 함께 HEART Lab 교육 커리큘럼을 개선',
], 16, DGRAY, 6)

# 3단계 카드
stages = [
    ('1단계', '문제 해결 사고 훈련', '"이 문제 어떻게 풀래?"\n가 바로 떠오를 때까지\n반복 훈련', NAVY),
    ('2단계', '모델을 도구로 쓰는 능력', '"YOLO 써봤습니다" X\n"이 문제에 왜 YOLO인지"\n설명할 수 있는 것', RGBColor(0x3B,0x82,0xF6)),
    ('3단계', '서비스를 만드는 능력', '챗봇 웹사이트\n실시간 객체 탐지\n12주면 혼자 구현 가능', GREEN),
]
for i, (num, title, desc, c) in enumerate(stages):
    left = Inches(0.6 + i * 4.2)
    rrect(s, left, Inches(2.8), Inches(3.9), Inches(3.2), c)
    txt(s, left, Inches(2.9), Inches(3.9), Inches(0.4),
        num, 15, RGBColor(0xBB,0xDD,0xFF), False, PP_ALIGN.CENTER)
    txt(s, left, Inches(3.3), Inches(3.9), Inches(0.5),
        title, 20, WHITE, True, PP_ALIGN.CENTER)
    txt(s, left + Inches(0.3), Inches(4.0), Inches(3.3), Inches(1.5),
        desc, 15, WHITE, False, PP_ALIGN.CENTER)

rrect(s, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.6), ORANGE)
txt(s, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.6),
    '목표:  문제를 던지면 해결할 수 있는 AI 개발자  |  부트캠프 6개월 → 12주 압축', 17, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: 학석사 혜택
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
slide_header(s, 'HEART Lab + 학석사')

# 왼쪽: 혜택 3개
# 시간
rrect(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(1.3), NAVY)
txt(s, Inches(0.9), Inches(1.55), Inches(2.5), Inches(0.5),
    '1년을 번다', 22, WHITE, True)
txt(s, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.4),
    '학부 4년 + 석사 2년 = 6년   →   학석사 = 5년', 15, RGBColor(0xBB,0xDD,0xFF))

# 돈
rrect(s, Inches(0.6), Inches(3.0), Inches(5.8), Inches(1.8), GREEN)
txt(s, Inches(0.9), Inches(3.1), Inches(5.2), Inches(0.5),
    '돈 받으면서 다닌다', 22, WHITE, True)
lines(s, Inches(0.9), Inches(3.7), Inches(5.2), Inches(1.0), [
    '학석사:  월 130만원',
    '석사 과정:  월 230만원',
    '*지원금 최대 지급',
], 16, WHITE, 4)

# 경력
rrect(s, Inches(0.6), Inches(5.0), Inches(5.8), Inches(1.3), ORANGE)
txt(s, Inches(0.9), Inches(5.1), Inches(5.2), Inches(0.5),
    '바로 취업 = 신입  →  학석사 = 주니어급', 18, WHITE, True)
txt(s, Inches(0.9), Inches(5.7), Inches(5.2), Inches(0.4),
    '시간 낭비가 아니라 가장 효율적인 투자', 14, WHITE)

# 오른쪽: CTA
rrect(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(4.8), RGBColor(0xE0,0xEB,0xFF))
txt(s, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.5),
    '궁금하면?', 26, NAVY, True, PP_ALIGN.CENTER)
lines(s, Inches(7.3), Inches(2.5), Inches(4.9), Inches(3.5), [
    '*연구실 견학 한번 와보세요',
    '실제로 뭘 하는지 보면 감이 옵니다',
    '',
    '*학부연구생으로 한 학기 해보고',
    '결정해도 됩니다',
    '',
    '',
    '(연락처 / QR코드)',
], 17, NAVY, 6)

# ============================================================
# SLIDE 12: 마무리
# ============================================================
s = prs.slides.add_slide(LY_BLANK)
rect(s, 0, 0, Inches(13.333), Inches(7.5), NAVY)

txt(s, Inches(1), Inches(2.2), Inches(11.3), Inches(0.8),
    '문제를 던지면 해결할 수 있는 사람이', 36, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.8),
    '살아남는다.', 44, WHITE, True, PP_ALIGN.CENTER)

rect(s, Inches(4.5), Inches(4.3), Inches(4.3), Inches(0.03), SKY)

txt(s, Inches(1), Inches(4.8), Inches(11.3), Inches(0.5),
    'HEART Lab  |  세종대학교', 20, SKY, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.4), Inches(11.3), Inches(0.4),
    'Human Emotion and Intelligent Agent Research for Future Transformation', 13, RGBColor(0x55,0x77,0xAA), False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.3),
    'github.com/RyanAhn533/HEART-Lab-Curriculum', 12, RGBColor(0x44,0x66,0x99), False, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
out = 'C:/Users/Ryan/HEART-Lab-Curriculum/딥러닝실습_발표_왜하는걸까.pptx'
prs.save(out)
print(f'OK: {out}')
