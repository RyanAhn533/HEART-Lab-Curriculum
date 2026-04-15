# -*- coding: utf-8 -*-
"""안준영 발표자료 수정 — 복사본에 작업"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_원본.pptx')

NV=RGBColor(0x01,0x3B,0x94)
DG=RGBColor(0x33,0x33,0x33)
GR=RGBColor(0x66,0x66,0x66)
SOFT=RGBColor(0x44,0x44,0x55)  # 연한 톤
FT='Noto Sans KR'

# === 1. 하트랩 연구 분야 추가 (S7, S8, S9) ===
# 빨간색 "하트랩 연구 분야 추가 :" 텍스트를 실제 연구 주제로 교체
research_line = "멀티모달 감정 AI  |  차량 운전자 감정 인식  |  협동로봇 AI  |  디지털 트윈  |  국제공동 (토론토대) Physical AI  |  문체부 버츄얼 생성형 AI"

for si in [6, 7, 8]:  # S7, S8, S9
    s = prs.slides[si]
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                has_red = False
                for r in p.runs:
                    try:
                        c = r.font.color.rgb
                        if c and str(c) == 'FF0000':
                            has_red = True
                            break
                    except:
                        pass
                if has_red:
                    # 빨간색 런들을 모두 지우고 새 텍스트로 교체
                    for r in p.runs:
                        try:
                            c = r.font.color.rgb
                            if c and str(c) == 'FF0000':
                                r.text = research_line
                                r.font.color.rgb = NV
                                r.font.size = Pt(11)
                                r.font.bold = False
                                # 나머지 빨간 런 비우기
                                found_first = True
                        except:
                            pass
                    # 깔끔하게: 빨간색이 아닌 나머지 런에서 "하트랩 연구 분야 추가" 관련 텍스트 제거
                    for r in p.runs:
                        if '연구 분야' in r.text or '추가' in r.text:
                            r.text = ''

# === 2. S9 국제공동 과제 내용 채우기 ===
s9 = prs.slides[8]
for sh in s9.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            # "국제공동 과제" 제목은 유지
            if p.text.strip() == '관련 내용':
                p.clear()
                # 첫 번째 "관련 내용" → 과제 설명
                break

# S9의 "관련 내용" 텍스트들을 찾아서 교체
for sh in s9.shapes:
    if sh.has_text_frame:
        tf = sh.text_frame
        for pi, p in enumerate(tf.paragraphs):
            t = p.text.strip()
            if t == '관련 내용':
                p.clear()

# S9에 국제공동 내용 추가 — 기존 텍스트박스 중 "관련 내용"이 있던 곳에
for sh in s9.shapes:
    if sh.has_text_frame:
        # 제목 "국제공동 과제" 옆의 상세 텍스트 박스 찾기
        if sh.left > 800000 and sh.top > 2500000 and sh.top < 3500000:
            tf = sh.text_frame
            if any('관련' in p.text or p.text.strip() == '' for p in tf.paragraphs):
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "국제공동 과제"
                p.font.size = Pt(28); p.font.bold = True; p.font.name = FT
                p.font.color.rgb = DG
                break

for sh in s9.shapes:
    if sh.has_text_frame:
        if sh.left > 800000 and sh.top > 3500000:
            tf = sh.text_frame
            if any('관련' in p.text or p.text.strip() == '' for p in tf.paragraphs):
                tf.clear()
                lines = [
                    "Project OptiBatt-AI",
                    "",
                    "성우하이텍(주관) + 토론토대학교(해외) + 세종대(AI) 컨소시엄",
                    "",
                    "EV 배터리 디지털 트윈 + AI Agent 기반 제조 자동화",
                    "멀티모달 센싱 데이터 → AI 해석 → 자동 설계-검증 루프",
                    "",
                    "세종대 역할:",
                    "  RGB-CT 데이터셋 구축 / 이미지-그래프 변환 모델 / 디지털 트윈 AI Agent",
                    "",
                    "토론토대학교:",
                    "  AI 노벨물리학상 2024 수상 대학",
                    "  생성형 AI, uPINN, 고속 Surrogate 등 AI 원천기술 파트너",
                ]
                for li, line in enumerate(lines):
                    if li == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    bold = line.startswith("Project") or line.startswith("토론토") or line.startswith("세종대 역할")
                    p.text = line
                    p.font.size = Pt(13)
                    p.font.name = FT
                    p.font.bold = bold
                    p.font.color.rgb = DG if not bold else NV
                    p.space_after = Pt(2)
                break

# === 3. S12 인턴/학사(학석사)/석사/박사 4개로 재구성 ===
s12 = prs.slides[11]
# 기존 텍스트 수정 — "학석사" 카드를 "학사"로 바꾸고 아래에 학석사연계 추가
for sh in s12.shapes:
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                # "학석사" → "학사" (메인 타이틀)
                if r.text.strip() == '학석사' and sh.left > 3000000 and sh.left < 5000000:
                    r.text = '학사'
                # 학석사 설명 부분에 연계프로그램 추가
                if '1년' in r.text and '빨리' in r.text:
                    r.text = '학석사 연계프로그램 가능'
                    r.font.size = Pt(14)
                if '잘하는 학생' in r.text:
                    r.text = '(1년 빨리 졸업, 잘하는 학생 한정)'
                    r.font.size = Pt(11)

# === 4. 전체 글씨 톤 연하게 (강조 부분, 흰색, 테마색 제외) ===
# 명시적으로 검정(rgb 설정된 것)만 연하게. 흰색/테마색/None은 절대 건드리지 않음
for si, s in enumerate(prs.slides):
    for sh in s.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        # color.type이 None이면 테마색 → 건드리지 않음
                        if r.font.color.type is None:
                            continue
                        c = r.font.color.rgb
                        if c is None:
                            continue
                        cs = str(c)
                        # 빨간색 이미 처리됨
                        if cs == 'FF0000':
                            continue
                        # 흰색 절대 건드리지 않음
                        if cs == 'FFFFFF' or cs == 'ffffff':
                            continue
                        # 밝은 색 (배경 위 흰색류) 건드리지 않음
                        cr = int(cs[:2], 16)
                        cg = int(cs[2:4], 16)
                        cb = int(cs[4:6], 16)
                        if cr > 180 and cg > 180 and cb > 180:
                            continue
                        # bold거나 큰 폰트는 건드리지 않음
                        if r.font.bold:
                            continue
                        sz = r.font.size
                        if sz and sz >= Pt(24):
                            continue
                        # 진한 검정(0~60)만 부드럽게
                        if cr < 60 and cg < 60 and cb < 60:
                            r.font.color.rgb = SOFT
                    except:
                        pass

out = 'C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_수정본.pptx'
prs.save(out)
print('OK:', out)
