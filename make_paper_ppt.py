# -*- coding: utf-8 -*-
"""S-PACE 논문 심사 PPT — S1~S3 유지 + 최신 논문 반영"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

# 원본 열어서 S1~S3 복사
orig = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/논문ppt_원본.pptx')
prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/논문ppt_원본.pptx')

# S4 이후 삭제 (S1~S3만 유지 = index 0,1,2)
while len(prs.slides) > 3:
    rId = prs.slides._sldIdLst[3].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[3])

# Colors & helpers
NV=RGBColor(0x01,0x3B,0x94); SK=RGBColor(0x86,0xB7,0xFE); W=RGBColor(0xFF,0xFF,0xFF)
DG=RGBColor(0x33,0x33,0x33); GR=RGBColor(0x66,0x66,0x66); LG=RGBColor(0xF0,0xF2,0xF5)
OR=RGBColor(0xFF,0x6B,0x35); GN=RGBColor(0x00,0xA8,0x6B); RD=RGBColor(0xE0,0x3E,0x3E)
LB=RGBColor(0xE0,0xEB,0xFF); FT='Noto Sans KR'; LY=prs.slide_layouts[0] # 1_제목 및 내용

FIGS='C:/Users/Ryan/HEART-Lab-Curriculum/paper_figs/'

def R(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def RR(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def T(s,l,t,w,h,text,sz=18,c=DG,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=FT;p.alignment=a
def ML(s,l,t,w,h,items,sz=16,c=DG,sp=6):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bold=item.startswith('*')
        if bold:item=item[1:]
        p.text=item;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=bold;p.font.name=FT;p.space_after=Pt(sp)
def PH(s,l,t,w,h,label):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=RGBColor(0xF8,0xF8,0xF8)
    sh.line.color.rgb=GR;sh.line.dash_style=2
    T(s,l,t+h//2-Inches(0.15),w,Inches(0.3),label,10,GR,False,PP_ALIGN.CENTER)
def IMG(s,path,l,t,width=None,height=None):
    w=width;h=height
    fp=FIGS+path if not os.path.isabs(path) else path
    if not os.path.exists(fp): PH(s,l,t,w or Inches(4),h or Inches(3),f'[{path}]'); return
    kw={};
    if w:kw['width']=w
    if h:kw['height']=h
    s.shapes.add_picture(fp,l,t,**kw)
def HDR(s,title,sub=None):
    # 템플릿 스타일 헤더
    R(s,0,0,Inches(13.333),Inches(1.0),SK)
    R(s,0,Inches(1.0),Inches(13.333),Inches(0.04),NV)
    T(s,Inches(0.5),Inches(0.15),Inches(12),Inches(0.6),title,28,NV,True)
    if sub:
        T(s,Inches(0.5),Inches(0.65),Inches(12),Inches(0.3),sub,13,GR)

# ============================================================
# S4. 핵심 아이디어 — Bio-Anchored + S-PACE
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'핵심 아이디어','Bio-Anchored Alignment + S-PACE')

RR(s,Inches(0.5),Inches(1.3),Inches(12.3),Inches(1.5),NV)
T(s,Inches(0.8),Inches(1.35),Inches(11.7),Inches(0.5),'Two Fundamental Challenges in Multimodal Emotion Recognition',20,W,True)
ML(s,Inches(0.8),Inches(1.9),Inches(11.5),Inches(0.8),
   ['*(1) Temporal Misalignment:  행동 반응(ms) vs 생리 반응(1~5초) — bio-behavioral lag',
    '*(2) Modality Noise:  얼굴 가림, 음성 왜곡, 센서 아티팩트 — 모달리티별 신뢰도 차이'],15,W,4)

RR(s,Inches(0.5),Inches(3.1),Inches(12.3),Inches(1.5),LB)
T(s,Inches(0.8),Inches(3.15),Inches(11.7),Inches(0.5),'Our Solution: SGMT (S-PACE Gated Multimodal Transformer)',20,NV,True)
ML(s,Inches(0.8),Inches(3.7),Inches(11.5),Inches(0.8),
   ['*S-PACE:  Bio를 Anchor(K,V)로, Behavioral(Video,Audio)을 Query로 — cross-attention 정렬',
    '*Quality-Aware Gate:  신호 품질 기반 동적 모달리티 가중치 (temperature annealing)',
    '*Bio-Anchored 설계:  생리신호 = 비자발적, 항상 존재 → 안정적 fusion 기준점'],14,NV,3)

T(s,Inches(0.5),Inches(5.0),Inches(12.3),Inches(0.4),
  '기존 MER: 모달리티를 대칭적으로 처리  →  SGMT: Bio를 중심에 놓고 Behavioral이 conditioning',16,GR)

PH(s,Inches(0.5),Inches(5.6),Inches(12.3),Inches(1.2),'[SGMT Architecture Figure — sgmt_architecture_v3.pdf]')

# ============================================================
# S5. 전체 아키텍처 (SGMT overview)
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'SGMT Architecture','End-to-end Pipeline')

PH(s,Inches(0.3),Inches(1.2),Inches(12.7),Inches(5.5),
   '[sgmt_architecture_v3.pdf — 전체 아키텍처 그림]\n\nVideo(ResNet) + Audio(AST) + Bio(ResNet per channel, GASF/GADF/MTF)\n→ Positional + Type + Speaker Embedding\n→ S-PACE Cross-Attention (Bio=K,V / Behavioral=Q)\n→ Quality-Aware Gate (temperature annealing)\n→ Temporal Swin Transformer (hierarchical)\n→ Perceiver Fusion (learnable latent tokens)\n→ Multi-task: Binary A/V + Quadrant + Polar')

# ============================================================
# S6. S-PACE 상세
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'S-PACE: Stimulus-Phase Alignment via Cross-modal Embedding')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(3.0),LB)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'설계 원리',20,NV,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(2.0),
   ['Behavioral(Video+Audio) = Query','Bio sensor channels = Key, Value','',
    '"현재 행동 관측에 대해 어떤','생리 채널이 가장 관련있는가?"','',
    '*행동이 먼저 발생 → 생리가 뒤따름','*→ Q로 행동을 쓰는 게 인과적으로 맞음'],14,NV,3)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(3.0),LG)
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'수식',18,DG,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(2.0),
   ['Q = W_Q [V_hat; A_hat]     (T x d)',
    'K = W_K B                   (C x d)',
    'V_att = W_V B               (C x d)',
    '',
    'B_t = MLP(softmax(QK^T/sqrt(d)) V_att)',
    '',
    '*→ Time-aligned bio representation'],13,DG,3)

RR(s,Inches(0.5),Inches(4.5),Inches(12.2),Inches(1.2),NV)
T(s,Inches(0.8),Inches(4.6),Inches(11.7),Inches(0.4),'Graceful Degradation',18,W,True)
ML(s,Inches(0.8),Inches(5.1),Inches(11.5),Inches(0.5),
   ['Video 없을 때 (KEMDy20): V_hat = 0 → Query가 Audio-only로 자동 전환',
    '*Bio anchor는 항상 존재 → 아키텍처 변경 없이 modality 가감 가능'],14,W,3)

PH(s,Inches(0.5),Inches(6.0),Inches(12.2),Inches(0.8),'[S-PACE attention heatmap — space_gate_heatmap.png]')

# ============================================================
# S7. Quality-Aware Gate
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Quality-Aware Gate','Dynamic Modality Fusion')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(2.5),LB)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'Quality Proxies',18,NV,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(1.5),
   ['q_face = 얼굴 검출 비율 (0~1)','q_audio = spectral flatness 기반','q_bio = 채널별 dynamic range 평균','','→ batch-wise z-score 정규화'],14,NV,4)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(2.5),LG)
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'Temperature Annealing',18,DG,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(1.5),
   ['w = softmax(MLP([V;A;B;q]) / tau)','','tau: 2.0 → 1.0 (cosine schedule)','','초기: 균등 탐색 → 후기: 전문화'],14,DG,4)

RR(s,Inches(0.5),Inches(4.0),Inches(12.2),Inches(1.0),NV)
T(s,Inches(0.8),Inches(4.05),Inches(11.7),Inches(0.4),'역할: interpretability module',18,W,True)
T(s,Inches(0.8),Inches(4.5),Inches(11.7),Inches(0.4),
  'Q-Acc 기여는 미미(-0.1%)하지만, 모달리티 중요도를 사람이 읽을 수 있게 해줌 (50K params, <0.05%)',14,W)

PH(s,Inches(0.5),Inches(5.3),Inches(12.2),Inches(1.5),'[Gate weight comparison — gate_comparison.png]')

# ============================================================
# S8. Temporal Swin + Perceiver Fusion
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Temporal Swin Transformer + Perceiver Fusion')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(3.5),LB)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'Hierarchical Temporal Swin',20,NV,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(2.5),
   ['1D Swin Transformer adaptation','','Stage 0: T=16→8 (window=4, shifted)',
    'Stage 1: T=8→4','Stage 2: T=4 (refinement)','',
    '4개 병렬 Encoder:','  1 post-gate (fused) + 3 intra-modal','  (video, audio, bio 각각 독립)'],14,NV,3)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(3.5),LG)
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'Perceiver Fusion',20,DG,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(2.5),
   ['Learnable latent tokens (L=32)','','Cross-attention:','  latent queries ← multimodal tokens','',
    '→ 고정 차원 emotion embedding','→ Multi-task heads:','   Binary (A/V), Quadrant, Polar'],14,DG,3)

T(s,Inches(0.5),Inches(5.0),Inches(12.2),Inches(0.4),
  'Multi-scale temporal + heterogeneous modality → compact representation → joint prediction',16,NV,True)

# ============================================================
# S9. Dataset
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Datasets','KEMDy20 + K-EmoCon')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(4.0),LB)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'KEMDy20',22,NV,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(3.0),
   ['40 sessions, 80 participants','Dyadic emotional conversations','',
    '*Modalities: Audio + Bio (EDA, TEMP, IBI)','*Video 없음 → skip_video=True','',
    'Arousal/Valence 5-point → binary','3-fold LOPPO CV','',
    '비교적 balanced distribution'],14,NV,4)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(4.0),RGBColor(0xFE,0xF0,0xE8))
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'K-EmoCon',22,OR,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(3.0),
   ['32 participants, 16 pairs','Paired conversation','',
    '*Modalities: Audio + Video + Bio','*(BVP, EDA, TEMP, ACC)','',
    'Valence: 5.1% Low / 94.9% High','→ extreme class imbalance','6-fold GroupKFold CV'],14,DG,4)

IMG(s,'class_distribution.png',Inches(0.5),Inches(5.5),width=Inches(12.2))

# ============================================================
# S10. Result — 메인 성능
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Results — Main Performance')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(3.0),NV)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'KEMDy20 (Audio + Bio)',20,W,True)
T(s,Inches(0.8),Inches(1.8),Inches(5.2),Inches(0.3),'3-fold LOPPO, skip-video + label smoothing',13,SK)
ML(s,Inches(0.8),Inches(2.3),Inches(5.0),Inches(1.5),
   ['*A-UAR: 68.4 ± 6.2%','*V-UAR: 57.6 ± 4.4%','*Q-Acc: 44.7 ± 4.8%  (chance=25%)','',
    'A-F1: 68.5%  |  V-F1: 55.4%'],18,W,5)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(3.0),RGBColor(0x3B,0x82,0xF6))
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'K-EmoCon (Audio + Video + Bio)',20,W,True)
T(s,Inches(7.1),Inches(1.8),Inches(5.3),Inches(0.3),'6-fold GroupKFold, 25 epochs',13,RGBColor(0xBB,0xDD,0xFF))
ML(s,Inches(7.1),Inches(2.3),Inches(5.3),Inches(1.5),
   ['*A-UAR: 62.9 ± 10.6%','*A-F1: 73.5%  |  V-F1: 86.1%','*Q-Acc: 62.5 ± 12.3%  (chance=25%)','',
    'V-UAR 미보고 (5.1% minority → 불안정)'],18,W,5)

RR(s,Inches(0.5),Inches(4.5),Inches(12.2),Inches(1.0),LG)
T(s,Inches(0.8),Inches(4.6),Inches(11.7),Inches(0.4),'Cross-Dataset Summary',18,DG,True)
ML(s,Inches(0.8),Inches(5.1),Inches(11.5),Inches(0.4),
   ['KEMDy20 (A+Bio): A-UAR 68.4%  |  K-EmoCon (A+V+Bio): A-UAR 62.9%',
    '*동일 아키텍처가 modality 가감에 자동 적응 (Gate mechanism)'],14,DG,3)

IMG(s,'curves_KEMDy20_3fold.png',Inches(0.5),Inches(5.8),width=Inches(6.0))
IMG(s,'confusion_KEMDy20_3fold_aggregated.png',Inches(6.8),Inches(5.8),width=Inches(6.0))

# ============================================================
# S11. Ablation
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Ablation Study')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(3.0),LB)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'KEMDy20 Architectural Ablation',16,NV,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(2.2),
   ['Audio-only:  A-UAR 71.4 / V-UAR 49.0','Bio-only:    A-UAR 54.2 / V-UAR 52.0',
    'Concat:      A-UAR 65.9 / V-UAR 51.3','*SGMT Full:  A-UAR 69.6 / V-UAR 52.8','',
    'Audio가 arousal 높지만 valence 최약','*→ SGMT = 가장 높은 composite score'],13,NV,3)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(3.0),RGBColor(0xFE,0xF0,0xE8))
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'K-EmoCon Component Ablation (6-fold)',16,OR,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(2.2),
   ['SGMT Full:      Q-Acc 63.2%','- Binary Head:   Q-Acc 59.0% (Δ-4.2)','- Polar Head:    Q-Acc 60.7% (Δ-2.5)',
    '- Gate Anneal:   Q-Acc 61.1% (Δ-2.1)','- Speaker Emb:   Q-Acc 61.8% (Δ-1.4)',
    '- Quality Gate:  Q-Acc 63.1% (Δ-0.1)','','*Binary head이 가장 critical (+4.2%)'],13,DG,3)

RR(s,Inches(0.5),Inches(4.5),Inches(12.2),Inches(1.8),LG)
T(s,Inches(0.8),Inches(4.6),Inches(11.7),Inches(0.4),'Key Findings',18,DG,True)
ML(s,Inches(0.8),Inches(5.1),Inches(11.5),Inches(1.0),
   ['*Bio contribution:  Bio 추가 시 V-UAR +3.8%p (audio-only 49.0 → SGMT 52.8) — valence에서 보완적',
    '*Gated > Concat:  A-UAR +3.7%p (65.9 → 69.6) — cross-attention + gating의 효과',
    '*SkipVideo:  video encoder 비활성화 시 +1.4% — zero-value noise 제거 효과'],13,DG,4)

# ============================================================
# S12. Gate Weight Analysis
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Gate Weight Analysis','Adaptive Modality Selection')

RR(s,Inches(0.5),Inches(1.2),Inches(5.8),Inches(2.5),NV)
T(s,Inches(0.8),Inches(1.3),Inches(5.2),Inches(0.5),'K-EmoCon (with Video)',18,W,True)
ML(s,Inches(0.8),Inches(1.9),Inches(5.0),Inches(1.5),
   ['*Video: ~95.2%','Audio: ~3.0%','Bio: ~1.8%','',
    '얼굴 표정이 가장 informative'],18,W,4)

RR(s,Inches(6.8),Inches(1.2),Inches(5.9),Inches(2.5),RGBColor(0x3B,0x82,0xF6))
T(s,Inches(7.1),Inches(1.3),Inches(5.3),Inches(0.5),'KEMDy20 (no Video)',18,W,True)
ML(s,Inches(7.1),Inches(1.9),Inches(5.3),Inches(1.5),
   ['Video: ~1.6% (zero tokens)','*Audio: ~89.7%','Bio: ~8.7%','',
    'Video 없으면 Audio로 자동 전환'],18,W,4)

IMG(s,'gate_comparison.png',Inches(0.5),Inches(4.0),width=Inches(6.0))
IMG(s,'space_gate_heatmap.png',Inches(6.8),Inches(4.0),width=Inches(6.0))

T(s,Inches(0.5),Inches(6.8),Inches(12.2),Inches(0.3),
  'Temporal dynamics: t=1에서 audio 57% + bio 13% → t≥4에서 video 96% (초기 acoustic context → steady-state visual tracking)',12,GR)

# ============================================================
# S13. Conclusion
# ============================================================
s=prs.slides.add_slide(LY)
HDR(s,'Conclusion & Future Work')

RR(s,Inches(0.5),Inches(1.2),Inches(12.2),Inches(3.0),NV)
T(s,Inches(0.8),Inches(1.3),Inches(11.7),Inches(0.5),'Contributions',22,W,True)
ML(s,Inches(0.8),Inches(1.9),Inches(11.5),Inches(2.2),
   ['*1. S-PACE:  bio-behavioral temporal lag를 cross-attention으로 정렬 (bio=anchor)',
    '*2. Quality-Aware Gate:  temperature annealing으로 modality 신뢰도 기반 동적 fusion',
    '*3. Unified Architecture:  modality 가감에 아키텍처 변경 없이 적응 (KEMDy20 ↔ K-EmoCon)',
    '*4. Gate Analysis:  supervision 없이 modality selection 학습 (video 95% ↔ audio 90%)'],16,W,6)

RR(s,Inches(0.5),Inches(4.5),Inches(5.8),Inches(2.2),LB)
T(s,Inches(0.8),Inches(4.6),Inches(5.2),Inches(0.5),'Limitations',18,NV,True)
ML(s,Inches(0.8),Inches(5.1),Inches(5.0),Inches(1.5),
   ['Bio gate weight ~2% → bio representation 개선 필요','K-EmoCon valence 94.9% skew → UAR 불안정',
    'Single-fold ablation의 일반화 한계'],14,NV,4)

RR(s,Inches(6.8),Inches(4.5),Inches(5.9),Inches(2.2),LG)
T(s,Inches(7.1),Inches(4.6),Inches(5.3),Inches(0.5),'Future Work',18,DG,True)
ML(s,Inches(7.1),Inches(5.1),Inches(5.3),Inches(1.5),
   ['NormWear (wearable foundation model)로','bio encoder 교체 → representation 개선',
    '','DEAP 데이터셋 일반화 검증','','Bio signal → LLM token space 직접 주입 (BioToken)'],14,DG,4)

# Save
out='C:/Users/Ryan/HEART-Lab-Curriculum/SPACE_논문심사_v1.pptx'
prs.save(out)
print('OK:',out)
