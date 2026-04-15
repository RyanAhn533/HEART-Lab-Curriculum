# -*- coding: utf-8 -*-
"""안준영 발표자료 수정 v3 — 원본에서 3가지만 수정"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_원본.pptx')

NV = RGBColor(0x01, 0x3B, 0x94)
W = RGBColor(0xFF, 0xFF, 0xFF)
DG = RGBColor(0x33, 0x33, 0x33)
FT = 'Noto Sans KR'

research_line = "멀티모달 감정 AI  |  차량 운전자 감정 인식  |  협동로봇 AI  |  디지털 트윈  |  국제공동 (토론토대) Physical AI  |  문체부 버츄얼 생성형 AI"

# ============================================================
# 1. S7, S8, S9: 빨간색 → 연구 주제 한 줄
# ============================================================
for si in [6, 7, 8]:
    s = prs.slides[si]
    for sh in s.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            red_runs = [r for r in p.runs if str(getattr(r.font.color, 'rgb', '')) == 'FF0000']
            if red_runs:
                red_runs[0].text = research_line
                red_runs[0].font.color.rgb = NV
                red_runs[0].font.size = Pt(9)
                red_runs[0].font.bold = False
                for r in red_runs[1:]:
                    r.text = ''

# ============================================================
# 2. S9: "관련 내용" → 국제공동 과제 (흰색 글씨)
# ============================================================
s9 = prs.slides[8]
for sh in s9.shapes:
    if not sh.has_text_frame: continue
    t_inch = round(sh.top / 914400, 2)
    full = sh.text_frame.text.strip()

    # "관련 내용" 첫번째 (T≈3.46) → 과제명
    if full == '관련 내용' and 3.0 < t_inch < 4.0:
        sh.text_frame.clear()
        p = sh.text_frame.paragraphs[0]
        p.text = "OptiBatt-AI  —  토론토대학교 × 세종대 × 성우하이텍"
        p.font.size = Pt(18); p.font.bold = True; p.font.name = FT; p.font.color.rgb = W

    # "관련 내용" 두번째 (T≈4.26) → 상세
    if full == '관련 내용' and t_inch > 4.0:
        sh.text_frame.clear()
        lines = [
            ("EV 배터리 디지털 트윈 + AI Agent 기반 제조 자동화", True),
            ("", False),
            ("토론토대학교 — AI 노벨물리학상 2024 수상 대학", True),
            ("  생성형 AI, uPINN, Surrogate 등 AI 원천기술 파트너", False),
            ("", False),
            ("세종대 역할:", True),
            ("  RGB-CT 데이터셋 / 이미지-그래프 변환 / 디지털 트윈 AI Agent", False),
        ]
        for li, (text, bold) in enumerate(lines):
            p = sh.text_frame.paragraphs[0] if li == 0 else sh.text_frame.add_paragraph()
            p.text = text; p.font.size = Pt(13); p.font.name = FT
            p.font.bold = bold; p.font.color.rgb = W; p.space_after = Pt(3)

# ============================================================
# 3. S12: 4개 카드 균등 배치 + 내용 수정 + 글씨색 수정
# ============================================================
s12 = prs.slides[11]

# 기존 shape 중 카드 관련(박스+텍스트+금액) 전부 삭제하고 새로 그림
# 유지할 것: 제목("AI 취업시장 현실"), 하단 CTA("궁금하면..."), 슬라이드번호
shapes_to_keep = []
for sh in s12.shapes:
    keep = False
    if sh.has_text_frame:
        full = sh.text_frame.text.strip()
        t_inch = round(sh.top / 914400, 2)
        # 제목 (T < 1.0)
        if t_inch < 1.0:
            keep = True
        # 하단 CTA (T > 6.0)
        if t_inch > 6.0:
            keep = True
        # 슬라이드 번호
        if full in ['11', '10', '12']:
            keep = True
    if keep:
        shapes_to_keep.append(sh)

# 삭제할 shape의 xml element 모으기
to_remove = []
for sh in s12.shapes:
    if sh not in shapes_to_keep:
        to_remove.append(sh._element)

for elem in to_remove:
    elem.getparent().remove(elem)

# 새로 4개 카드 그리기
def add_rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def add_txt(slide, l, t, w, h, text, sz=16, color=DG, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = FT; p.alignment = align

def add_ml(slide, l, t, w, h, items, sz=13, color=DG):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = FT; p.space_after = Pt(3)

# 4카드 균등: 전체폭 12.7", 카드폭 2.9", 간격 0.2", 시작 0.3"
card_w = Inches(2.9)
gap = Inches(0.2)
start_x = Inches(0.35)
card_h = Inches(4.3)
card_top = Inches(1.5)

cards = [
    {
        'title': '인턴',
        'color': RGBColor(0xF0, 0xF2, 0xF5),
        'txt_color': DG,
        'money': '월 20만원 + a',
        'lines': [
            ('3개월 체험 후 선택 가능', False),
            ('', False),
            ('인턴 이후', False),
            ('학석사 연계 가능', True),
        ],
    },
    {
        'title': '학사',
        'color': NV,
        'txt_color': W,
        'money': '월 130만원 + a',
        'lines': [
            ('학석사 연계프로그램 가능', True),
            ('(1년 빨리 졸업,', False),
            ('잘하는 학생에 한해서)', False),
            ('', False),
            ('바로 취업 = 신입', False),
            ('학석사 = 주니어급', True),
        ],
    },
    {
        'title': '석사',
        'color': RGBColor(0x00, 0xA8, 0x6B),
        'txt_color': W,
        'money': '월 230만원 + a',
        'lines': [
            ('학부 졸업 후 진학', False),
            ('', False),
            ('지원금 최대 지급', True),
        ],
    },
    {
        'title': '박사',
        'color': RGBColor(0xF0, 0xF2, 0xF5),
        'txt_color': DG,
        'money': '월 320만원',
        'lines': [
            ('석사 졸업 후 진학', False),
            ('', False),
            ('지원금 최대 지급', True),
        ],
    },
]

for i, card in enumerate(cards):
    x = int(start_x + i * (card_w + gap))

    # 배경 박스
    add_rect(s12, x, card_top, card_w, card_h, card['color'])

    # 타이틀
    add_txt(s12, x, card_top + Inches(0.2), card_w, Inches(0.5),
            card['title'], 26, card['txt_color'], True)

    # 설명
    add_ml(s12, x + Inches(0.2), card_top + Inches(1.0), card_w - Inches(0.4), Inches(2.5),
           card['lines'], 14, card['txt_color'])

    # 금액
    add_txt(s12, x, card_top + Inches(3.5), card_w, Inches(0.5),
            card['money'], 22, card['txt_color'], True)

out = 'C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_수정본.pptx'
prs.save(out)
print('OK:', out)
