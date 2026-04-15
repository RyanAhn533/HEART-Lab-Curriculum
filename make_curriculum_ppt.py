"""
HEART Lab AI Curriculum — PPT 생성 스크립트
기존 안준영_발표자료_수정본.pptx 템플릿 기반
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# 기존 PPT를 템플릿으로 사용 (슬라이드 마스터/레이아웃 가져오기)
template_path = os.path.join(os.path.dirname(__file__), '안준영_발표자료_수정본.pptx')
prs = Presentation(template_path)

# 기존 슬라이드 전부 삭제
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# 레이아웃
layout_title = prs.slide_layouts[1]   # 제목 슬라이드
layout_content = prs.slide_layouts[0] # 1_제목 및 내용
layout_blank = prs.slide_layouts[7]   # 빈 화면

# ── 색상 ──
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x2D, 0x5B, 0xE3)
RED = RGBColor(0xCC, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── 헬퍼 ──
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = '맑은 고딕'
    p.alignment = alignment
    return tf

def add_bullet_list(tf, items, font_size=13, color=DARK, indent=0):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = '맑은 고딕'
        p.level = indent
        p.space_before = Pt(3)

def add_slide_number(slide, num):
    add_textbox(slide, 0.5, 0.5, 1.5, 1, str(num), font_size=24, bold=True, color=ACCENT)

# ══════════════════════════════════════
# Slide 1: 타이틀
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, 3, 5, 27, 4,
    'HEART Lab AI Curriculum',
    font_size=36, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 3, 9, 27, 2,
    '신입 연구원 AI 교육 과정',
    font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 3, 12, 27, 1.5,
    '안준영  |  세종대학교 HEART Lab  |  2026',
    font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# Slide 2: 교육 이념
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 1)
add_textbox(slide, 2, 0.5, 28, 2,
    '교육 이념',
    font_size=28, bold=True, color=DARK)
add_textbox(slide, 2, 3, 28, 2,
    '"AI로 문제를 풀 수 있는 사람을 만든다"',
    font_size=22, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, 2, 6, 13, 8,
    '시장 현실', font_size=16, bold=True, color=RED)
add_bullet_list(tf, [
    'AI 경력직 요구: 54% → 80.6%',
    '2026년 신입 채용: 12.4%',
    '기업은 바로 투입 가능한 사람만 원함',
], font_size=12)

tf = add_textbox(slide, 17, 6, 13, 8,
    '바로 투입이란', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '"YOLO 써봤습니다" X',
    '"왜 YOLO인지 설명 + 설계 + 구현" O',
    '코딩은 기본, 뭘 만들지 아는 것까지',
], font_size=12)

# ══════════════════════════════════════
# Slide 3: Top-Down Thinking
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 2)
add_textbox(slide, 2, 0.5, 28, 2,
    'Top-Down Thinking',
    font_size=28, bold=True, color=DARK)
add_textbox(slide, 2, 3, 28, 1.5,
    'First Principles — 근본으로 돌아가서 문제를 재정의',
    font_size=16, color=GRAY)

tf = add_textbox(slide, 2, 5.5, 13, 8,
    '연구에서', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '기존 DMS: 졸음 감지, 주시 이탈',
    '→ "왜 졸음만 보지?"',
    '→ 근본 원인 = 감정 변화 누적',
    '→ Frustration 기반 재정의',
    '→ 현대차와 협업 이유',
], font_size=12)

tf = add_textbox(slide, 17, 5.5, 13, 8,
    '교육에서', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '수학부터 X',
    '일단 돌려보고 결과 확인',
    '→ "왜 이렇게 되지?" 파고듦',
    '→ 필요한 이론을 역추적',
    '→ 이해한 후 다시 코드로',
], font_size=12)

# ══════════════════════════════════════
# Slide 4: 3단계 성장 모델
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 3)
add_textbox(slide, 2, 0.5, 28, 2,
    '3단계 성장 모델',
    font_size=28, bold=True, color=DARK)

# 1단계
tf = add_textbox(slide, 2, 3.5, 9, 7,
    '1단계', font_size=18, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '문제 해결 사고',
    '',
    '"이 문제 어떻게 풀래?"',
    '바로 떠오를 때까지',
    '반복 훈련',
    '',
    'Phase 2~3',
], font_size=11)

# 2단계
tf = add_textbox(slide, 12, 3.5, 9, 7,
    '2단계', font_size=18, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '모델을 도구로',
    '',
    '"이 문제에 왜 이 모델?"',
    '설명할 수 있는 것',
    '',
    'Phase 4~8',
], font_size=11)

# 3단계
tf = add_textbox(slide, 22, 3.5, 9, 7,
    '3단계', font_size=18, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '서비스 구현',
    '',
    '챗봇, 실시간 탐지',
    '혼자 만들 수 있는 것',
    '',
    'Phase 9~10',
], font_size=11)

# ══════════════════════════════════════
# Slide 5: 코드 5대 원칙
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 4)
add_textbox(slide, 2, 0.5, 28, 2,
    '코드 5대 원칙',
    font_size=28, bold=True, color=DARK)

principles = [
    ('구조 통일', '모든 코드가 #0→#1→#2→#3→#4'),
    ('하나만 바뀐다', '이전 코드에서 변경은 항상 1개 섹션'),
    ('같은 데이터, 다른 기법', '5개 데이터셋 끝까지 반복'),
    ('간결하게', 'TF2/Keras 간결한 API만'),
    ('60~120줄', '파일당 최대 120줄'),
]
for i, (title, desc) in enumerate(principles):
    y = 3.5 + i * 2.2
    add_textbox(slide, 3, y, 5, 1.5, title, font_size=16, bold=True, color=ACCENT)
    add_textbox(slide, 9, y, 20, 1.5, desc, font_size=13, color=DARK)

# ══════════════════════════════════════
# Slide 6: 코드 뼈대
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 5)
add_textbox(slide, 2, 0.5, 28, 2,
    '코드 뼈대 — 모든 파일이 이 구조',
    font_size=28, bold=True, color=DARK)

code_text = (
    '#0. 라이브러리\n'
    '    from tensorflow.keras.models import Sequential\n'
    '    from sklearn.model_selection import train_test_split\n\n'
    '#1. 데이터\n'
    '    x, y = load_boston().data, load_boston().target\n\n'
    '#2. 데이터 전처리 및 분할\n'
    '    x_train, x_test = train_test_split(...)\n'
    '    scaler.fit_transform(x_train)\n\n'
    '#3. 모델\n'
    '    model = Sequential()\n'
    '    model.add(Dense(64, activation=\'relu\'))\n'
    '    model.compile(loss=\'mse\', optimizer=\'adam\')\n'
    '    model.fit(x_train, y_train, epochs=100)\n\n'
    '#4. 평가\n'
    '    loss = model.evaluate(x_test, y_test)\n'
    '    y_pred = model.predict(x_test)'
)
add_textbox(slide, 3, 3, 26, 13, code_text, font_size=11, color=DARK)

# ══════════════════════════════════════
# Slide 7: 전체 로드맵
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 6)
add_textbox(slide, 2, 0.5, 28, 2,
    '전체 로드맵 — Phase 0~10',
    font_size=28, bold=True, color=DARK)

roadmap_left = [
    ('Phase 0', '왜 하는가', GRAY),
    ('Phase 1', '통계/수학 기초', GRAY),
    ('Phase 2', '돌려봐', ACCENT),
    ('Phase 3', '왜 안 되지?', ACCENT),
    ('Phase 4', '저장, 구조', ACCENT),
    ('Phase 5', '이미지 (CNN)', ACCENT),
]
roadmap_right = [
    ('Phase 6', '시계열 (RNN)', ACCENT),
    ('Phase 7', 'NLP', ACCENT),
    ('Phase 8', '성능 올리기', ACCENT),
    ('Phase 9', 'ML 기본기', ACCENT),
    ('Phase 10', 'PyTorch', ACCENT),
]

for i, (phase, desc, color) in enumerate(roadmap_left):
    y = 3 + i * 2
    add_textbox(slide, 2, y, 5, 1.5, phase, font_size=13, bold=True, color=color)
    add_textbox(slide, 7, y, 8, 1.5, desc, font_size=12, color=DARK)

for i, (phase, desc, color) in enumerate(roadmap_right):
    y = 3 + i * 2
    add_textbox(slide, 17, y, 5, 1.5, phase, font_size=13, bold=True, color=color)
    add_textbox(slide, 22, y, 8, 1.5, desc, font_size=12, color=DARK)

add_textbox(slide, 2, 15.5, 28, 1,
    'TF2/Keras (Phase 0~8)  →  sklearn (Phase 9)  →  PyTorch (Phase 10)',
    font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# Slide 8: "하나만 바뀐다" 예시
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 7)
add_textbox(slide, 2, 0.5, 28, 2,
    '"하나만 바뀐다" — 학습의 핵심',
    font_size=28, bold=True, color=DARK)

changes = [
    ('02 기본 회귀', '#0~#4 뼈대 자체를 익힘', '—'),
    ('03 다중입력', '#3에 레이어 추가', '#3'),
    ('04 train/test', '#2에 split 추가', '#2'),
    ('06 EarlyStopping', '#3에 콜백 추가', '#3'),
    ('07 이진분류', '#3에 sigmoid + loss 변경', '#3'),
    ('09 Scaler', '#2에 StandardScaler 추가', '#2'),
    ('13 CNN', '#1 이미지, #2 reshape, #3 Conv2D', '#1,#2,#3'),
    ('17 LSTM', '#3에 RNN → LSTM 교체', '#3'),
]
for i, (step, desc, section) in enumerate(changes):
    y = 3 + i * 1.7
    add_textbox(slide, 2, y, 6, 1.2, step, font_size=11, bold=True, color=ACCENT)
    add_textbox(slide, 8.5, y, 14, 1.2, desc, font_size=11, color=DARK)
    add_textbox(slide, 24, y, 5, 1.2, section, font_size=11, bold=True, color=RED)

# ══════════════════════════════════════
# Slide 9: 이론 학습 원칙
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 8)
add_textbox(slide, 2, 0.5, 28, 2,
    '이론 학습 원칙',
    font_size=28, bold=True, color=DARK)
add_textbox(slide, 2, 3, 28, 1.5,
    '영상 (직관)  →  블로그 리뷰 (정리)  →  논문 (원본)',
    font_size=18, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 2, 5, 28, 1,
    '절대 논문부터 읽지 마라.',
    font_size=14, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, 2, 7, 9, 8,
    '1순위: 3Blue1Brown', font_size=15, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '시각적 직관 최강',
    '한글자막 있음',
    '이해 안 되면 여기부터',
    '',
    'Neural Networks',
    'Linear Algebra',
    'Calculus',
    'Transformers',
], font_size=10)

tf = add_textbox(slide, 12, 7, 9, 8,
    '2순위: 구글 검색', font_size=15, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '"[주제] 논문 리뷰" 검색',
    '',
    '예시:',
    '"LSTM 논문 리뷰"',
    '"ResNet 논문 리뷰"',
    '"Word2Vec 논문 리뷰"',
    '',
    '한국어 블로그 2~3개',
], font_size=10)

tf = add_textbox(slide, 22, 7, 9, 8,
    '3순위: 추천 채널', font_size=15, bold=True, color=ACCENT)
add_bullet_list(tf, [
    'codingopera',
    '  (Transformer 정리)',
    '',
    'StatQuest',
    '  (통계/ML, 영어)',
    '',
    '혁펜하임',
    '  (한국어 DL)',
], font_size=10)

# ══════════════════════════════════════
# Slide 10: HEART Lab 연결
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 9)
add_textbox(slide, 2, 0.5, 28, 2,
    '배우는 것이 어디에 쓰이는가',
    font_size=28, bold=True, color=DARK)

tf = add_textbox(slide, 2, 3.5, 9, 9,
    '현대자동차', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '차량 감정인식',
    '멀티모달 센싱',
    '→ 운전자 감정 → 안전운전',
    '',
    'Phase 5~7',
    '(CNN, RNN, NLP)',
], font_size=11)

tf = add_textbox(slide, 12, 3.5, 9, 9,
    '두산', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '협동로봇 AI',
    'VLM/VLA + AI SoC',
    '+ 강화학습',
    '',
    'Phase 8~10',
    '(전이학습, PyTorch)',
], font_size=11)

tf = add_textbox(slide, 22, 3.5, 9, 9,
    '토론토대', font_size=16, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '국제공동 과제',
    'EV 배터리 디지털 트윈',
    '+ AI Agent',
    '',
    'Phase 10',
    '(논문 재현)',
], font_size=11)

# ══════════════════════════════════════
# Slide 11: 운영 방식
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_slide_number(slide, 10)
add_textbox(slide, 2, 0.5, 28, 2,
    '운영 방식',
    font_size=28, bold=True, color=DARK)

tf = add_textbox(slide, 2, 3, 13, 10,
    '단일 트랙', font_size=18, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '학부/석사 구분 없음',
    '전부 같은 트랙',
    '아는 사람은 스킵',
    '모르면 처음부터',
    '',
    '스킵 기준:',
    '30분 안에 혼자 짤 수 있으면 넘어감',
], font_size=12)

tf = add_textbox(slide, 17, 3, 13, 10,
    '평가', font_size=18, bold=True, color=ACCENT)
add_bullet_list(tf, [
    '체크포인트 1: Phase 3 끝',
    '  (회귀/분류 체화)',
    '',
    '체크포인트 2: Phase 6 끝',
    '  (CNN/RNN 체화)',
    '',
    '최종 발표: Phase 10 끝',
    '  (논문 재현 프로젝트)',
    '',
    '매월 월간 발표',
], font_size=12)

# ══════════════════════════════════════
# Slide 12: 한 문장 요약
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, 3, 4, 27, 10,
    '#0→#1→#2→#3→#4 뼈대를 몸에 익히고,\n\n'
    '한 번에 하나씩만 바꿔가며,\n\n'
    '같은 데이터에 다른 기법을 적용하면서,\n\n'
    '3Blue1Brown과 논문 리뷰로 "왜"를 이해하고,\n\n'
    'HEART Lab의 실제 과제에 투입될 수 있는\n'
    '연구원을 만든다.',
    font_size=18, color=DARK, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════
# Slide 13: 감사합니다
# ══════════════════════════════════════
slide = prs.slides.add_slide(layout_blank)
add_textbox(slide, 3, 7, 27, 3,
    '감사합니다',
    font_size=36, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

# ── 저장 ──
output_path = os.path.join(os.path.dirname(__file__), 'CURRICULUM_MASTER.pptx')
prs.save(output_path)
print(f'저장 완료: {output_path}')
