# -*- coding: utf-8 -*-
"""v5 — CLAUDE.md 적용. 한 장에 한 메시지, 여백 충분히, 빠진 내용 복원."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/template.pptx')
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Colors
NV=RGBColor(0x01,0x3B,0x94); SK=RGBColor(0x86,0xB7,0xFE); W=RGBColor(0xFF,0xFF,0xFF)
DG=RGBColor(0x33,0x33,0x33); GR=RGBColor(0x66,0x66,0x66); LG=RGBColor(0xF0,0xF2,0xF5)
OR=RGBColor(0xFF,0x6B,0x35); GN=RGBColor(0x00,0xA8,0x6B); RD=RGBColor(0xE0,0x3E,0x3E)
LB=RGBColor(0xE0,0xEB,0xFF); FT='Noto Sans KR'; LY=prs.slide_layouts[7]

def R(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def RR(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def T(s,l,t,w,h,text,sz=18,c=DG,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=FT;p.alignment=a
def ML(s,l,t,w,h,items,sz=16,c=DG,sp=8):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bold=item.startswith('*')
        if bold:item=item[1:]
        p.text=item;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=bold;p.font.name=FT;p.space_after=Pt(sp)
def PH(s,l,t,w,h,label):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=RGBColor(0xF8,0xF8,0xF8)
    sh.line.color.rgb=GR;sh.line.dash_style=2
    T(s,l,t+h//2-Inches(0.15),w,Inches(0.3),label,10,GR,False,PP_ALIGN.CENTER)
def HDR(s,title):
    R(s,0,0,Inches(13.333),Inches(1.05),SK);R(s,0,Inches(1.05),Inches(13.333),Inches(0.04),NV)
    T(s,Inches(0.6),Inches(0.15),Inches(12),Inches(0.7),title,30,NV,True)
    R(s,0,Inches(7.2),Inches(13.333),Inches(0.3),NV);T(s,Inches(0.4),Inches(7.22),Inches(5),Inches(0.25),'HEART Lab | Sejong Univ.',9,W)

# ============================================================
# S1. 타이틀
# ============================================================
s=prs.slides.add_slide(LY)
R(s,0,0,Inches(13.333),Inches(7.5),SK)
R(s,0,Inches(2.5),Inches(13.333),Inches(2.8),NV)
R(s,0,Inches(7.1),Inches(13.333),Inches(0.4),NV)
T(s,Inches(0.8),Inches(0.7),Inches(5),Inches(0.5),'2026-1학기 딥러닝실습',16,NV)
T(s,Inches(1),Inches(2.9),Inches(11.3),Inches(1.0),'딥러닝 실습, 왜 하는 걸까?',48,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(4.0),Inches(11.3),Inches(0.5),'대학, 취업, 그리고 AI',22,RGBColor(0xBB,0xD5,0xFF),False,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(6.2),Inches(11.3),Inches(0.5),'안준영  |  세종대학교 HEART Lab',18,NV,False,PP_ALIGN.CENTER)

# ============================================================
# S2. 취업 현실 — 기사 헤드라인으로 때리기
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'AI 취업시장,  지금 어떤 상황인가')

# 기사 헤드라인 placeholder 3개 (캡처해서 넣을 자리)
PH(s,Inches(0.6),Inches(1.4),Inches(6.0),Inches(1.5),
   '[기사] "인턴조차 3년차 경력자" — 경향신문 2026.03')
PH(s,Inches(0.6),Inches(3.1),Inches(6.0),Inches(1.5),
   '[기사] "코딩만으론 어림없다…신입 취업 대재앙" — 글로벌이코노믹 2025.12')
PH(s,Inches(0.6),Inches(4.8),Inches(6.0),Inches(1.5),
   '[기사] "AI 도입 기업, 신입 채용 40% 감소" — 경향신문 2025.09')

# 통계
RR(s,Inches(7.0),Inches(1.4),Inches(5.7),Inches(1.3),RD)
T(s,Inches(7.3),Inches(1.5),Inches(5.1),Inches(0.4),'AI 경력직 요구 비율',15,W)
T(s,Inches(7.3),Inches(1.95),Inches(5.1),Inches(0.5),'54%  →  80.6%',32,W,True)

RR(s,Inches(7.0),Inches(2.9),Inches(5.7),Inches(1.3),RD)
T(s,Inches(7.3),Inches(3.0),Inches(5.1),Inches(0.4),'2026 신입 채용 비율',15,W)
T(s,Inches(7.3),Inches(3.45),Inches(5.1),Inches(0.5),'12.4%',32,W,True)

RR(s,Inches(7.0),Inches(4.4),Inches(5.7),Inches(1.3),RD)
T(s,Inches(7.3),Inches(4.5),Inches(5.1),Inches(0.4),'SW 신입 비중 (2년 변화)',15,W)
T(s,Inches(7.3),Inches(4.95),Inches(5.1),Inches(0.5),'53.5%  →  37.4%',32,W,True)

T(s,Inches(7.0),Inches(5.9),Inches(5.7),Inches(0.3),'출처: 경향신문, ZDNet, 서울신문, 한국은행',10,GR)

RR(s,Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.5),NV)
T(s,Inches(0.6),Inches(6.4),Inches(12.1),Inches(0.5),
  '학부 졸업만으로는 AI 취업이 안 되는 시대',20,W,True,PP_ALIGN.CENTER)

# ============================================================
# S3. AI 취업 두 갈래
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'인공지능 취업,  크게 두 가지')

RR(s,Inches(0.6),Inches(1.4),Inches(5.8),Inches(5.0),LB)
T(s,Inches(0.6),Inches(1.6),Inches(5.8),Inches(0.6),'기업 개발자',28,NV,True,PP_ALIGN.CENTER)
ML(s,Inches(1.2),Inches(2.5),Inches(4.8),Inches(2.5),
   ['각 회사마다 도메인이 다름','(반도체, 자동차, 의학, 금융, 로봇...)','',
    '그 도메인에서','*서비스(제품)를 만드는 사람'],18,DG,8)
PH(s,Inches(1.5),Inches(5.5),Inches(1.0),Inches(0.5),'[삼성]')
PH(s,Inches(2.8),Inches(5.5),Inches(1.0),Inches(0.5),'[현대]')
PH(s,Inches(4.1),Inches(5.5),Inches(1.0),Inches(0.5),'[네이버]')

RR(s,Inches(6.9),Inches(1.4),Inches(5.8),Inches(5.0),RGBColor(0xFE,0xF0,0xE8))
T(s,Inches(6.9),Inches(1.6),Inches(5.8),Inches(0.6),'연구소 (연구원)',28,OR,True,PP_ALIGN.CENTER)
ML(s,Inches(7.5),Inches(2.5),Inches(4.8),Inches(2.5),
   ['기업 연구소, 국책 연구소, 대학원','','트렌드에 맞춰 기술 발전을 리드','논문, 특허, 기술 이전','',
    '*연구소는 국책과제로 운영됨'],18,DG,8)

RR(s,Inches(2.0),Inches(6.6),Inches(9.3),Inches(0.4),NV)
T(s,Inches(2.0),Inches(6.6),Inches(9.3),Inches(0.4),
  '공통:  AI로 문제를 풀 수 있는 사람이 필요',18,W,True,PP_ALIGN.CENTER)

# ============================================================
# S4. 왜 안뽑냐 → 바로 투입
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'왜 신입을 안 뽑는가')

T(s,Inches(1),Inches(2.0),Inches(11.3),Inches(0.8),
  '교육시켜놔도  1~2년 만에 이직',30,DG,True,PP_ALIGN.CENTER)

T(s,Inches(1),Inches(3.2),Inches(11.3),Inches(0.5),
  '그래서 처음부터',22,GR,False,PP_ALIGN.CENTER)

RR(s,Inches(2.5),Inches(4.2),Inches(8.3),Inches(1.5),NV)
T(s,Inches(2.5),Inches(4.5),Inches(8.3),Inches(0.8),
  '바로 투입 가능한 사람',36,W,True,PP_ALIGN.CENTER)

T(s,Inches(1),Inches(6.2),Inches(11.3),Inches(0.4),
  '스펙이 아니라,  문제를 풀어본 경험이 있는 사람',20,NV,False,PP_ALIGN.CENTER)

# ============================================================
# S5. 문제 해결 능력
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'바로 투입  =  문제 해결 능력')

RR(s,Inches(0.6),Inches(1.4),Inches(12.1),Inches(1.8),LB)
T(s,Inches(1.0),Inches(1.5),Inches(11.3),Inches(0.5),
  '기업은  "YOLO 돌려주세요"  라고 안 함',22,NV,True)
T(s,Inches(1.0),Inches(2.2),Inches(11.3),Inches(0.5),
  '"우리 공장에서 불량품 찾고 싶은데"  →  여기서부터 다 너가 설계',20,DG)

RR(s,Inches(0.6),Inches(3.5),Inches(5.5),Inches(3.0),LG)
T(s,Inches(0.9),Inches(3.7),Inches(5.1),Inches(0.5),'코딩은 당연히 해야 한다',20,DG,True)
ML(s,Inches(0.9),Inches(4.4),Inches(4.8),Inches(1.8),
   ['코드를 짜고 핸들링하는 건 기본','','*근데 코딩만으로는 부족','*뭘 만들지 + 어떻게 접근할지','*까지 되어야 함'],16,GR,6)

RR(s,Inches(6.5),Inches(3.5),Inches(6.2),Inches(3.0),NV)
T(s,Inches(6.8),Inches(3.7),Inches(5.6),Inches(0.5),'지금 핫한 기술',20,W,True)
ML(s,Inches(6.8),Inches(4.4),Inches(5.3),Inches(1.8),
   ['LLM  (ChatGPT, Claude)','AI Agent  (자율적 의사결정)','멀티모달  (영상 + 음성 + 텍스트)','Physical AI  (로봇, 자율주행)','',
    '*이걸 조합해서 구현 = 바로 투입'],16,W,5)

# ============================================================
# S6. 세계 트렌드 — 감정 AI
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'세계는 지금 어디로 가고 있냐')

T(s,Inches(0.6),Inches(1.3),Inches(12),Inches(0.5),
  '감정 AI (Emotion AI)  —  글로벌 트렌드',24,DG,True)

cards=[
    ('NVIDIA','Audio2Face:\n감정 분석 → 표정 생성\n\nAuto Emotion\n\nGTC 2026\nAffective Care Drive',NV),
    ('Affectiva\n(Smart Eye)','감정 AI 글로벌 1위\n\n차량 감정인식 상용화\n\n$73.5M 인수',RGBColor(0x3B,0x82,0xF6)),
    ('Top Conference','CVPR, AAAI,\nINTERSPEECH\n\n감정 인식이\n가장 활발한 주제',GN),
]
for i,(title,desc,c) in enumerate(cards):
    left=Inches(0.6+i*4.2)
    RR(s,left,Inches(2.0),Inches(3.9),Inches(3.8),c)
    T(s,left,Inches(2.2),Inches(3.9),Inches(0.7),title,22,W,True,PP_ALIGN.CENTER)
    T(s,left,Inches(3.1),Inches(3.9),Inches(2.4),desc,15,W,False,PP_ALIGN.CENTER)

PH(s,Inches(0.8),Inches(6.1),Inches(3.5),Inches(0.7),'[CES2024 NVIDIA 감정기술]')
PH(s,Inches(4.8),Inches(6.1),Inches(3.5),Inches(0.7),'[GTC2026 Block Diagram]')
PH(s,Inches(8.8),Inches(6.1),Inches(3.5),Inches(0.7),'[Affectiva In-Cabin]')

# ============================================================
# S7. 현대차
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'우리는 이미 거기에 있다  ①')

RR(s,Inches(0.6),Inches(1.3),Inches(12.1),Inches(0.7),LB)
T(s,Inches(0.8),Inches(1.35),Inches(11.7),Inches(0.5),
  'HEART Lab  —  Human Emotion and Intelligent Agent Research for Future Transformation',15,NV,True,PP_ALIGN.CENTER)

RR(s,Inches(0.6),Inches(2.3),Inches(12.1),Inches(4.2),NV)
PH(s,Inches(9.0),Inches(2.5),Inches(3.5),Inches(1.0),'[현대자동차 로고]')
T(s,Inches(1.0),Inches(2.5),Inches(7.5),Inches(0.7),'현대자동차',32,W,True)
T(s,Inches(1.0),Inches(3.2),Inches(7.5),Inches(0.5),'미래기술 공모전  —  선정',22,SK,True)

ML(s,Inches(1.0),Inches(4.0),Inches(11.0),Inches(2.0),
   ['기존에 하던 감정인식 연구를 발전시켜 제안  →  바로 뽑힘','',
    'Affective Care Drive:  멀티모달 센싱 → 감정 인식 → 안전 운전 지원','',
    '*이미 하고 있던 기술을 조금 바꿔서 낸 건데 바로 됨',
    '*=  기술력이 이미 그 수준이었다는 뜻'],17,W,5)

# ============================================================
# S8. 두산
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'우리는 이미 거기에 있다  ②')

RR(s,Inches(0.6),Inches(1.3),Inches(12.1),Inches(5.2),RGBColor(0x3B,0x82,0xF6))
PH(s,Inches(9.0),Inches(1.5),Inches(3.5),Inches(1.0),'[두산 로보틱스 로고]')
T(s,Inches(1.0),Inches(1.5),Inches(7.5),Inches(0.7),'두산 로보틱스',32,W,True)
T(s,Inches(1.0),Inches(2.2),Inches(7.5),Inches(0.5),'역대 가장 큰 금액 규모의 R&D 과제',22,RGBColor(0xBB,0xDD,0xFF),True)

ML(s,Inches(1.0),Inches(3.0),Inches(11.0),Inches(2.5),
   ['협동로봇에 AI를 넣는 프로젝트','',
    'VLM / VLA  +  AI SoC  +  강화학습','',
    'AI 모델 경량화,  멀티모달 데이터,  인간-로봇 협업(HRI) 담당','',
    '*두산이 먼저 찾아온 것'],18,W,7)

RR(s,Inches(2.5),Inches(6.6),Inches(8.3),Inches(0.4),OR)
T(s,Inches(2.5),Inches(6.6),Inches(8.3),Inches(0.4),
  '대기업에서 먼저 찾아오는 연구실',20,W,True,PP_ALIGN.CENTER)

# ============================================================
# S9. First Principles
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'어떻게 이렇게 됐냐')

# 머스크
RR(s,Inches(0.6),Inches(1.3),Inches(5.8),Inches(2.8),LG)
PH(s,Inches(0.8),Inches(1.5),Inches(2.2),Inches(1.2),'[일론 머스크]')
T(s,Inches(3.2),Inches(1.5),Inches(3.0),Inches(0.5),'First Principles',20,DG,True)
T(s,Inches(3.2),Inches(2.0),Inches(3.0),Inches(0.4),'Thinking',20,DG,True)
ML(s,Inches(0.8),Inches(2.9),Inches(5.4),Inches(1.0),
   ['"로켓이 왜 비싸지?"  →  원자재 = 2%','*기존 상식 의심 → 근본으로 → 새로운 답'],14,DG,4)

# HEART Lab
RR(s,Inches(6.8),Inches(1.3),Inches(5.9),Inches(2.8),LB)
T(s,Inches(7.1),Inches(1.5),Inches(5.3),Inches(0.5),'HEART Lab',22,NV,True)
ML(s,Inches(7.1),Inches(2.2),Inches(5.3),Inches(1.5),
   ['기존 DMS = 졸음, 주시 이탈','→ 업계가 다 이렇게 함','',
    '*"잠깐,  왜 졸음만 보지?"'],16,NV,5)

# Frustration
RR(s,Inches(0.6),Inches(4.4),Inches(12.1),Inches(2.3),NV)
T(s,Inches(0.9),Inches(4.5),Inches(11.5),Inches(0.5),'근본으로 돌아감',24,W,True)
ML(s,Inches(0.9),Inches(5.1),Inches(11.0),Inches(1.4),
   ['운전 위험의 근본 원인  =  감정 변화가 누적되는 과정',
    '*→  Frustration(좌절감) 기반으로 문제를 재정의',
    '*→  현대차가 뽑은 이유가 바로 이것'],18,W,6)

# ============================================================
# S10. 커리큘럼
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'너도 할 수 있다')

ML(s,Inches(0.6),Inches(1.3),Inches(12),Inches(0.7),
   ['카카오 부트캠프 + 비트캠프 경험,  교수님과 함께 커리큘럼 개선',
    '*목표:  문제를 던지면 해결할 수 있는 AI 개발자'],15,DG,3)

stages=[
    ('1단계','문제 해결 사고','"이 문제 어떻게 풀래?"\n\n바로 떠오를 때까지',NV),
    ('2단계','모델을 도구로','"왜 이 모델인지"\n\n설명할 수 있는 것',RGBColor(0x3B,0x82,0xF6)),
    ('3단계','서비스 구현','챗봇, YOLO 등\n\n12주면 혼자 가능',GN),
]
for i,(num,title,desc,c) in enumerate(stages):
    left=Inches(0.6+i*4.2)
    RR(s,left,Inches(2.3),Inches(3.9),Inches(3.5),c)
    T(s,left,Inches(2.5),Inches(3.9),Inches(0.3),num,14,RGBColor(0xBB,0xDD,0xFF),False,PP_ALIGN.CENTER)
    T(s,left,Inches(2.9),Inches(3.9),Inches(0.5),title,22,W,True,PP_ALIGN.CENTER)
    T(s,left+Inches(0.3),Inches(3.6),Inches(3.3),Inches(1.8),desc,16,W,False,PP_ALIGN.CENTER)

T(s,Inches(0.6),Inches(6.1),Inches(12.1),Inches(0.4),
  '부트캠프 6개월  →  12주 압축',16,GR,False,PP_ALIGN.CENTER)

# ============================================================
# S11. 인턴 → 학석사
# ============================================================
s=prs.slides.add_slide(LY);HDR(s,'어떻게 시작하냐')

RR(s,Inches(0.6),Inches(1.4),Inches(3.8),Inches(5.0),LG)
T(s,Inches(0.9),Inches(1.7),Inches(3.2),Inches(0.5),'인턴',28,DG,True,PP_ALIGN.CENTER)
T(s,Inches(0.9),Inches(2.5),Inches(3.2),Inches(0.6),'월 20만원',24,NV,True,PP_ALIGN.CENTER)
ML(s,Inches(0.9),Inches(3.4),Inches(3.2),Inches(2.5),
   ['부담 없이 3개월 체험','','안 맞으면','그만둬도 됨'],18,DG,6)

RR(s,Inches(4.8),Inches(1.4),Inches(3.8),Inches(5.0),NV)
T(s,Inches(5.1),Inches(1.7),Inches(3.2),Inches(0.5),'학석사',28,W,True,PP_ALIGN.CENTER)
T(s,Inches(5.1),Inches(2.5),Inches(3.2),Inches(0.6),'월 130만원',24,W,True,PP_ALIGN.CENTER)
ML(s,Inches(5.1),Inches(3.4),Inches(3.2),Inches(2.5),
   ['1년 빨리 졸업 가능','(잘하는 학생에 한해서)','','바로 취업 = 신입','학석사 = 주니어급'],16,W,5)

RR(s,Inches(8.9),Inches(1.4),Inches(3.8),Inches(5.0),GN)
T(s,Inches(9.2),Inches(1.7),Inches(3.2),Inches(0.5),'석사',28,W,True,PP_ALIGN.CENTER)
T(s,Inches(9.2),Inches(2.5),Inches(3.2),Inches(0.6),'월 230만원',24,W,True,PP_ALIGN.CENTER)
ML(s,Inches(9.2),Inches(3.4),Inches(3.2),Inches(2.5),
   ['학부 졸업 후 진학','','*지원금 최대 지급'],18,W,6)

T(s,Inches(0.6),Inches(6.6),Inches(6.5),Inches(0.3),
  '궁금하면:  견학 오세요  |  인턴으로 해보고 결정',16,NV,True)
T(s,Inches(9.5),Inches(6.6),Inches(3.3),Inches(0.3),'(연락처 / QR코드)',14,GR,False,PP_ALIGN.RIGHT)

# ============================================================
# S12. 마무리
# ============================================================
s=prs.slides.add_slide(LY)
R(s,0,0,Inches(13.333),Inches(7.5),NV)
T(s,Inches(1),Inches(2.0),Inches(11.3),Inches(0.8),
  '문제를 던지면',38,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(2.9),Inches(11.3),Inches(0.8),
  '해결할 수 있는 사람이',38,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(4.0),Inches(11.3),Inches(1.0),
  '살아남는다.',50,W,True,PP_ALIGN.CENTER)
R(s,Inches(4.5),Inches(5.2),Inches(4.3),Inches(0.03),SK)
T(s,Inches(1),Inches(5.6),Inches(11.3),Inches(0.5),
  'HEART Lab  |  세종대학교',20,SK,False,PP_ALIGN.CENTER)

out='C:/Users/Ryan/HEART-Lab-Curriculum/딥러닝실습_발표_v5.pptx'
prs.save(out)
print('OK:',out)
