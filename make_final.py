# -*- coding: utf-8 -*-
"""딥러닝실습 발표 PPT 최종 — 정보 분산, 한 장에 한 메시지"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/template.pptx')
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

NV=RGBColor(0x01,0x3B,0x94); SK=RGBColor(0x86,0xB7,0xFE)
W=RGBColor(0xFF,0xFF,0xFF); DG=RGBColor(0x33,0x33,0x33)
GR=RGBColor(0x66,0x66,0x66); LG=RGBColor(0xF0,0xF2,0xF5)
OR=RGBColor(0xFF,0x6B,0x35); GN=RGBColor(0x00,0xA8,0x6B)
RD=RGBColor(0xE0,0x3E,0x3E); LB=RGBColor(0xE0,0xEB,0xFF)
FT='Noto Sans KR'; LY=prs.slide_layouts[7]

def R(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def RR(s,l,t,w,h,c):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=c;sh.line.fill.background()
def T(s,l,t,w,h,text,sz=18,c=DG,b=False,a=PP_ALIGN.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name=FT;p.alignment=a
def ML(s,l,t,w,h,items,sz=16,c=DG,sp=8):
    tb=s.shapes.add_textbox(l,t,w,h);tf=tb.text_frame;tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bold=item.startswith('*')
        if bold:item=item[1:]
        p.text=item;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=bold;p.font.name=FT;p.space_after=Pt(sp)
def PH(s,l,t,w,h,label):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h);sh.fill.solid();sh.fill.fore_color.rgb=RGBColor(0xF8,0xF8,0xF8)
    sh.line.color.rgb=GR;sh.line.dash_style=2
    T(s,l,t+h//2-Inches(0.15),w,Inches(0.3),label,10,GR,False,PP_ALIGN.CENTER)
def HDR(s,title):
    R(s,0,0,Inches(13.333),Inches(1.05),SK);R(s,0,Inches(1.05),Inches(13.333),Inches(0.04),NV)
    T(s,Inches(0.6),Inches(0.2),Inches(12),Inches(0.65),title,28,NV,True)
    R(s,0,Inches(7.2),Inches(13.333),Inches(0.3),NV);T(s,Inches(0.4),Inches(7.2),Inches(5),Inches(0.3),'HEART Lab | Sejong Univ.',9,W)

# ===== S1: 타이틀 =====
s=prs.slides.add_slide(LY)
R(s,0,0,Inches(13.333),Inches(7.5),SK)
R(s,0,Inches(2.2),Inches(13.333),Inches(3.2),NV)
R(s,0,Inches(7.1),Inches(13.333),Inches(0.4),NV)
T(s,Inches(0.8),Inches(0.6),Inches(5),Inches(0.5),'2026-1학기 딥러닝실습',16,NV)
T(s,Inches(1),Inches(2.7),Inches(11.3),Inches(1.0),'딥러닝 실습, 왜 하는 걸까?',44,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(3.8),Inches(11.3),Inches(0.5),'대학, 취업, 그리고 AI',20,RGBColor(0xBB,0xD5,0xFF),False,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(6.0),Inches(11.3),Inches(0.5),'안준영  |  세종대학교 HEART Lab',18,NV,False,PP_ALIGN.CENTER)

# ===== S2: 취업 현실 =====
s=prs.slides.add_slide(LY);HDR(s,'AI 취업시장 현실')

PH(s,Inches(0.6),Inches(1.4),Inches(6.0),Inches(5.2),'[사람인 AI 채용공고 캡처]\n\n경력 3년 이상만 뽑는 현실')

RR(s,Inches(7.0),Inches(1.4),Inches(5.7),Inches(1.4),RD)
T(s,Inches(7.3),Inches(1.5),Inches(5.1),Inches(0.5),'AI 경력직 요구 비율',16,W)
T(s,Inches(7.3),Inches(2.0),Inches(5.1),Inches(0.5),'54%  →  80.6%',30,W,True)

RR(s,Inches(7.0),Inches(3.0),Inches(5.7),Inches(1.4),RD)
T(s,Inches(7.3),Inches(3.1),Inches(5.1),Inches(0.5),'2026년 신입 채용 비율',16,W)
T(s,Inches(7.3),Inches(3.6),Inches(5.1),Inches(0.5),'12.4%',30,W,True)

RR(s,Inches(7.0),Inches(4.6),Inches(5.7),Inches(1.4),RD)
T(s,Inches(7.3),Inches(4.7),Inches(5.1),Inches(0.5),'SW 신입 비중 2년 변화',16,W)
T(s,Inches(7.3),Inches(5.2),Inches(5.1),Inches(0.5),'53.5%  →  37.4%',30,W,True)

T(s,Inches(7.0),Inches(6.2),Inches(5.7),Inches(0.3),'출처: ZDNet 2025.12 / 서울신문 2026.01',10,GR)

# ===== S3: AI 취업 두 갈래 =====
s=prs.slides.add_slide(LY);HDR(s,'인공지능 취업, 크게 두 가지')

RR(s,Inches(0.6),Inches(1.4),Inches(5.8),Inches(4.5),LB)
T(s,Inches(0.6),Inches(1.5),Inches(5.8),Inches(0.6),'기업 개발자',26,NV,True,PP_ALIGN.CENTER)
ML(s,Inches(1.0),Inches(2.3),Inches(5.0),Inches(3.0),
   ['삼성, 현대, LG, 네이버, 카카오...',
    '',
    '각 회사의 도메인이 다름',
    '(반도체, 자동차, 의학, 금융, 로봇...)',
    '',
    '*도메인에 맞는 서비스(제품)를',
    '*만드는 사람'],17,DG,5)
PH(s,Inches(1.5),Inches(5.2),Inches(1.2),Inches(0.5),'[삼성]')
PH(s,Inches(3.0),Inches(5.2),Inches(1.2),Inches(0.5),'[네이버]')
PH(s,Inches(4.5),Inches(5.2),Inches(1.2),Inches(0.5),'[카카오]')

RR(s,Inches(6.9),Inches(1.4),Inches(5.8),Inches(4.5),RGBColor(0xFE,0xF0,0xE8))
T(s,Inches(6.9),Inches(1.5),Inches(5.8),Inches(0.6),'연구소 (연구원)',26,OR,True,PP_ALIGN.CENTER)
ML(s,Inches(7.3),Inches(2.3),Inches(5.0),Inches(3.0),
   ['기업 연구소, 국책 연구소, 대학원',
    '',
    '트렌드에 맞춰 기술 발전을 리드',
    '논문, 특허, 기술 이전',
    '',
    '*연구소는 국책과제(정부 과제)로 운영'],17,DG,5)

RR(s,Inches(2.0),Inches(6.2),Inches(9.3),Inches(0.7),NV)
T(s,Inches(2.0),Inches(6.25),Inches(9.3),Inches(0.6),
  '공통:  AI로 문제를 풀 수 있는 사람이 필요하다',20,W,True,PP_ALIGN.CENTER)

# ===== S4: 왜 안뽑냐 =====
s=prs.slides.add_slide(LY);HDR(s,'왜 신입을 안 뽑는가')

T(s,Inches(0.8),Inches(1.8),Inches(11.5),Inches(0.6),
  '교육시켜놔도  1~2년 만에 이직해버림',28,DG,True,PP_ALIGN.CENTER)

T(s,Inches(0.8),Inches(3.0),Inches(11.5),Inches(0.6),
  '그래서 처음부터',24,GR,False,PP_ALIGN.CENTER)

RR(s,Inches(2.0),Inches(4.0),Inches(9.3),Inches(1.5),NV)
T(s,Inches(2.0),Inches(4.2),Inches(9.3),Inches(1.0),
  '바로 투입 가능한 사람을 원한다',32,W,True,PP_ALIGN.CENTER)

T(s,Inches(0.8),Inches(6.0),Inches(11.5),Inches(0.5),
  '스펙이 아니라,  문제를 풀어본 경험이 있는 사람',20,NV,False,PP_ALIGN.CENTER)

# ===== S4: 문제 해결 능력 =====
s=prs.slides.add_slide(LY);HDR(s,'바로 투입 = 문제 해결 능력')

RR(s,Inches(0.6),Inches(1.4),Inches(12.1),Inches(2.0),LB)
T(s,Inches(0.9),Inches(1.5),Inches(11.5),Inches(0.5),
  '기업은 "YOLO 돌려주세요" 라고 안 함',22,NV,True)
T(s,Inches(0.9),Inches(2.2),Inches(11.5),Inches(0.8),
  '"우리 공장에서 불량품 찾고 싶은데"  →  여기서부터 다 설계해야 함',20,DG)

RR(s,Inches(0.6),Inches(3.8),Inches(5.5),Inches(2.5),LG)
T(s,Inches(0.9),Inches(3.9),Inches(5.1),Inches(0.5),'코딩은 당연히 해야 한다',20,DG,True)
ML(s,Inches(0.9),Inches(4.5),Inches(4.8),Inches(1.5),
   ['코드를 짜고 핸들링하는 건 기본','','*근데 코딩만 할 줄 아는 건 부족','*뭘 만들지 아는 것까지 되어야 함'],16,GR,6)

RR(s,Inches(6.5),Inches(3.8),Inches(6.2),Inches(2.5),NV)
T(s,Inches(6.8),Inches(3.9),Inches(5.6),Inches(0.5),'지금 핫한 기술',20,W,True)
ML(s,Inches(6.8),Inches(4.5),Inches(5.3),Inches(1.5),
   ['LLM  (ChatGPT, Claude)','AI Agent  (자율적 의사결정)','멀티모달  (영상 + 음성 + 텍스트)','Physical AI  (로봇, 자율주행)'],16,W,6)

T(s,Inches(0.6),Inches(6.6),Inches(12.1),Inches(0.4),
  '이걸 조합해서 실제 서비스를 구현할 수 있는 능력  =  바로 투입',18,NV,True,PP_ALIGN.CENTER)

# ===== S5: 세계 트렌드 =====
s=prs.slides.add_slide(LY);HDR(s,'세계는 지금 어디로 가고 있냐')

T(s,Inches(0.6),Inches(1.3),Inches(12),Inches(0.5),
  '감정 AI (Emotion AI) — 글로벌 트렌드',24,DG,True)

cards=[
    ('NVIDIA','Audio2Face:\n감정 분석 → 표정 생성\n\nAuto Emotion\n\nGTC 2026\nAffective Care Drive',NV),
    ('Affectiva\n(Smart Eye)','감정 AI 글로벌 1위\n\n차량 감정인식 상용화\n\n$73.5M 인수',RGBColor(0x3B,0x82,0xF6)),
    ('Top Conference','CVPR, AAAI,\nINTERSPEECH\n\n감정 인식이\n가장 활발한 주제',GN),
]
for i,(title,desc,c) in enumerate(cards):
    left=Inches(0.6+i*4.2)
    RR(s,left,Inches(2.0),Inches(3.9),Inches(4.2),c)
    T(s,left,Inches(2.2),Inches(3.9),Inches(0.6),title,22,W,True,PP_ALIGN.CENTER)
    T(s,left,Inches(3.0),Inches(3.9),Inches(2.8),desc,15,W,False,PP_ALIGN.CENTER)

PH(s,Inches(0.8),Inches(6.5),Inches(3.5),Inches(0.5),'[CES2024 NVIDIA]')
PH(s,Inches(4.8),Inches(6.5),Inches(3.5),Inches(0.5),'[GTC2026]')
PH(s,Inches(8.8),Inches(6.5),Inches(3.5),Inches(0.5),'[Affectiva]')

# ===== S6: 현대차 =====
s=prs.slides.add_slide(LY);HDR(s,'우리는 이미 거기에 있다  ①')

RR(s,Inches(0.6),Inches(1.4),Inches(12.1),Inches(0.8),LB)
T(s,Inches(0.8),Inches(1.5),Inches(11.7),Inches(0.5),
  'HEART Lab  —  Human Emotion and Intelligent Agent Research for Future Transformation',16,NV,True,PP_ALIGN.CENTER)

RR(s,Inches(0.6),Inches(2.5),Inches(12.1),Inches(4.0),NV)
T(s,Inches(1.0),Inches(2.7),Inches(11.3),Inches(0.6),'현대자동차',30,W,True)
T(s,Inches(1.0),Inches(3.3),Inches(11.3),Inches(0.5),'미래기술 공모전  —  선정',22,SK,True)

ML(s,Inches(1.0),Inches(4.1),Inches(10.5),Inches(2.0),
   ['기존에 하던 감정인식 연구를 발전시켜 제안',
    '',
    'Affective Care Drive:',
    '멀티모달 센싱  →  운전자 감정 인식  →  안전 운전 지원',
    '',
    '*하고 있던 기술을 조금 바꿔서 낸 건데 바로 된 것',
    '*= 기술력이 이미 그 수준이었다는 뜻'],17,W,5)

PH(s,Inches(8.5),Inches(2.7),Inches(4.0),Inches(1.0),'[현대차 로고]')

# ===== S7: 두산 =====
s=prs.slides.add_slide(LY);HDR(s,'우리는 이미 거기에 있다  ②')

RR(s,Inches(0.6),Inches(1.4),Inches(12.1),Inches(4.8),RGBColor(0x3B,0x82,0xF6))
T(s,Inches(1.0),Inches(1.6),Inches(11.3),Inches(0.6),'두산 로보틱스',30,W,True)
T(s,Inches(1.0),Inches(2.2),Inches(11.3),Inches(0.5),'역대 가장 큰 금액 규모의 R&D 과제',22,RGBColor(0xBB,0xDD,0xFF),True)

ML(s,Inches(1.0),Inches(3.0),Inches(10.5),Inches(2.5),
   ['협동로봇에 AI를 넣는 프로젝트',
    '',
    'VLM / VLA  +  AI SoC  +  강화학습',
    '',
    'AI 모델 경량화,  멀티모달 데이터,  인간-로봇 협업(HRI) 담당',
    '',
    '*두산이 먼저 찾아온 것'],18,W,6)

PH(s,Inches(8.5),Inches(1.6),Inches(4.0),Inches(1.0),'[두산 로보틱스 로고]')

RR(s,Inches(2.0),Inches(6.5),Inches(9.3),Inches(0.5),OR)
T(s,Inches(2.0),Inches(6.5),Inches(9.3),Inches(0.5),
  '대기업에서 먼저 찾아오는 연구실',20,W,True,PP_ALIGN.CENTER)

# ===== S8: First Principles =====
s=prs.slides.add_slide(LY);HDR(s,'어떻게 이렇게 됐냐  —  First Principles')

# 머스크
RR(s,Inches(0.6),Inches(1.4),Inches(5.8),Inches(2.8),LG)
PH(s,Inches(0.8),Inches(1.5),Inches(2.2),Inches(1.2),'[일론 머스크]')
T(s,Inches(3.2),Inches(1.5),Inches(3.0),Inches(0.4),'First Principles',18,DG,True)
T(s,Inches(3.2),Inches(1.9),Inches(3.0),Inches(0.4),'Thinking',18,DG,True)
T(s,Inches(3.2),Inches(2.4),Inches(3.0),Inches(0.3),'— Elon Musk',13,GR)
ML(s,Inches(0.8),Inches(2.9),Inches(5.4),Inches(1.0),
   ['"로켓이 왜 비싸지?"  →  원자재 = 2%','→  SpaceX 10배 절감'],15,DG,4)

# HEART Lab
RR(s,Inches(6.8),Inches(1.4),Inches(5.9),Inches(2.8),LB)
T(s,Inches(7.1),Inches(1.5),Inches(5.3),Inches(0.5),'HEART Lab',20,NV,True)
ML(s,Inches(7.1),Inches(2.1),Inches(5.3),Inches(1.8),
   ['기존 DMS = 졸음 감지, 주시 이탈','→ 업계가 다 이렇게 함','',
    '*"잠깐, 왜 졸음만 보지?"'],16,NV,5)

# 아래: Frustration
RR(s,Inches(0.6),Inches(4.5),Inches(12.1),Inches(2.2),NV)
T(s,Inches(0.9),Inches(4.6),Inches(11.5),Inches(0.5),'근본으로 돌아감',22,W,True)
ML(s,Inches(0.9),Inches(5.2),Inches(11.0),Inches(1.3),
   ['운전 위험의 근본 원인  =  감정 변화가 누적되는 과정',
    '',
    '*Frustration(좌절감) 기반으로 문제를 재정의',
    '*→  기존에 아무도 안 하던 접근  →  현대차가 뽑은 이유가 이것'],17,W,5)

# ===== S9: 커리큘럼 =====
s=prs.slides.add_slide(LY);HDR(s,'너도 할 수 있다  —  커리큘럼')

ML(s,Inches(0.6),Inches(1.3),Inches(12),Inches(0.8),
   ['카카오 AI 부트캠프, 비트캠프 경험  +  교수님과 함께 교육 커리큘럼 개선',
    '*목표:  문제를 던지면 해결할 수 있는 AI 개발자'],16,DG,4)

stages=[
    ('1단계','문제 해결 사고','"이 문제 어떻게 풀래?"\n\n바로 떠오를 때까지\n반복 훈련',NV),
    ('2단계','모델을 도구로','"YOLO 써봤습니다" X\n\n"이 문제에 왜 YOLO인지"\n설명할 수 있는 것',RGBColor(0x3B,0x82,0xF6)),
    ('3단계','서비스 구현','챗봇 웹사이트\n실시간 객체 탐지\n\n12주면 혼자 가능',GN),
]
for i,(num,title,desc,c) in enumerate(stages):
    left=Inches(0.6+i*4.2)
    RR(s,left,Inches(2.4),Inches(3.9),Inches(3.8),c)
    T(s,left,Inches(2.5),Inches(3.9),Inches(0.4),num,14,RGBColor(0xBB,0xDD,0xFF),False,PP_ALIGN.CENTER)
    T(s,left,Inches(2.9),Inches(3.9),Inches(0.5),title,22,W,True,PP_ALIGN.CENTER)
    T(s,left+Inches(0.3),Inches(3.6),Inches(3.3),Inches(2.2),desc,16,W,False,PP_ALIGN.CENTER)

T(s,Inches(0.6),Inches(6.5),Inches(12.1),Inches(0.4),
  '부트캠프 6개월  →  12주 압축',16,GR,False,PP_ALIGN.CENTER)

# ===== S10: 인턴/학석사 =====
s=prs.slides.add_slide(LY);HDR(s,'어떻게 시작하냐')

RR(s,Inches(0.6),Inches(1.4),Inches(3.8),Inches(4.5),LG)
T(s,Inches(0.9),Inches(1.6),Inches(3.2),Inches(0.5),'인턴',26,DG,True,PP_ALIGN.CENTER)
ML(s,Inches(0.9),Inches(2.3),Inches(3.2),Inches(2.5),
   ['*월 20만원','','부담 없이 3개월 체험','','안 맞으면','그만둬도 됨'],18,DG,4)

RR(s,Inches(4.8),Inches(1.4),Inches(3.8),Inches(4.5),NV)
T(s,Inches(5.1),Inches(1.6),Inches(3.2),Inches(0.5),'학석사',26,W,True,PP_ALIGN.CENTER)
ML(s,Inches(5.1),Inches(2.3),Inches(3.2),Inches(2.5),
   ['*월 130만원','','1년 빨리 졸업 가능','(잘하는 학생에 한해서)','','바로 취업 = 신입','학석사 = 주니어급'],16,W,4)

RR(s,Inches(8.8),Inches(1.4),Inches(3.9),Inches(4.5),GN)
T(s,Inches(9.1),Inches(1.6),Inches(3.3),Inches(0.5),'석사',26,W,True,PP_ALIGN.CENTER)
ML(s,Inches(9.1),Inches(2.3),Inches(3.3),Inches(2.5),
   ['*월 230만원','','학부 졸업 후 진학','','*지원금 최대 지급'],18,W,4)

T(s,Inches(0.6),Inches(6.2),Inches(6),Inches(0.4),
  '궁금하면:  견학 오세요  |  인턴으로 해보고 결정',18,NV,True)
T(s,Inches(9.0),Inches(6.2),Inches(3.7),Inches(0.4),'(연락처 / QR코드)',16,GR,False,PP_ALIGN.RIGHT)

# ===== S11: 마무리 =====
s=prs.slides.add_slide(LY)
R(s,0,0,Inches(13.333),Inches(7.5),NV)
T(s,Inches(1),Inches(2.0),Inches(11.3),Inches(0.8),
  '문제를 던지면',36,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(2.8),Inches(11.3),Inches(0.8),
  '해결할 수 있는 사람이',36,W,True,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(3.8),Inches(11.3),Inches(1.0),
  '살아남는다.',48,W,True,PP_ALIGN.CENTER)
R(s,Inches(4.5),Inches(5.0),Inches(4.3),Inches(0.03),SK)
T(s,Inches(1),Inches(5.5),Inches(11.3),Inches(0.5),
  'HEART Lab  |  세종대학교',20,SK,False,PP_ALIGN.CENTER)
T(s,Inches(1),Inches(6.1),Inches(11.3),Inches(0.4),
  'Human Emotion and Intelligent Agent Research for Future Transformation',13,RGBColor(0x55,0x77,0xAA),False,PP_ALIGN.CENTER)

out='C:/Users/Ryan/HEART-Lab-Curriculum/딥러닝실습_발표_최종.pptx'
prs.save(out)
print('OK:',out)
