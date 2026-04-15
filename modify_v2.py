# -*- coding: utf-8 -*-
"""안준영 발표자료 수정 v2 — 딱 3가지만 수정, 나머지 일체 안 건드림"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation('C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_원본.pptx')

NV = RGBColor(0x01, 0x3B, 0x94)
research_line = "멀티모달 감정 AI  |  차량 운전자 감정 인식  |  협동로봇 AI  |  디지털 트윈  |  국제공동 (토론토대) Physical AI  |  문체부 버츄얼 생성형 AI"

# === 1. S7, S8, S9: 빨간색 텍스트 → 연구 주제 한 줄 ===
for si in [6, 7, 8]:
    s = prs.slides[si]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            first_red_done = False
            for r in p.runs:
                try:
                    c = r.font.color.rgb
                    if c and str(c) == 'FF0000':
                        if not first_red_done:
                            r.text = research_line
                            r.font.color.rgb = NV
                            r.font.size = Pt(10)
                            r.font.bold = False
                            first_red_done = True
                        else:
                            r.text = ''
                except:
                    pass

# === 2. S9: "관련 내용" → 국제공동 과제 내용 ===
s9 = prs.slides[8]
for sh in s9.shapes:
    if not sh.has_text_frame:
        continue
    tf = sh.text_frame
    # "국제공동 과제" 제목 텍스트박스는 유지
    # "관련 내용"이 있는 텍스트박스만 수정
    full = tf.text.strip()
    if full == '국제공동 과제':
        continue
    if '관련 내용' in full:
        # 이 텍스트박스의 위치로 어떤 박스인지 판단
        if sh.top > 2500000 and sh.top < 3500000:
            # 제목 아래 첫번째 "관련 내용"
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "Project OptiBatt-AI  —  토론토대학교 × 세종대 × 성우하이텍"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.name = 'Noto Sans KR'
            p.font.color.rgb = NV
        elif sh.top > 3500000:
            # 아래쪽 "관련 내용" → 상세 내용
            tf.clear()
            lines = [
                ("EV 배터리 디지털 트윈 + AI Agent 기반 제조 자동화", True, NV),
                ("", False, None),
                ("토론토대학교 — AI 노벨물리학상 2024 수상", True, NV),
                ("생성형 AI, uPINN, 고속 Surrogate 등 AI 원천기술 파트너", False, None),
                ("", False, None),
                ("세종대 역할:", True, NV),
                ("RGB-CT 데이터셋 구축 / 이미지-그래프 변환 모델", False, None),
                ("디지털 트윈 AI Agent 개발 / 멀티모달 센싱 데이터 해석", False, None),
            ]
            for li, (text, bold, color) in enumerate(lines):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(14)
                p.font.name = 'Noto Sans KR'
                p.font.bold = bold
                if color:
                    p.font.color.rgb = color
                p.space_after = Pt(3)

# === 3. S12: 인턴/학사(학석사연계)/석사/박사 4개 ===
s12 = prs.slides[11]
for sh in s12.shapes:
    if not sh.has_text_frame:
        continue
    tf = sh.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            t = r.text.strip()
            # "학석사" 메인 타이틀 → "학사"로 변경 (위치로 판별: 두번째 카드)
            if t == '학석사' and sh.left > 3000000 and sh.left < 6000000:
                r.text = '학사'
            # 학석사 설명 수정
            if '1년' in t and '빨리' in t:
                r.text = '학석사 연계프로그램 가능'
                r.font.size = Pt(14)
            if '잘하는 학생' in t:
                r.text = '(1년 빨리 졸업, 잘하는 학생 한정)'
                r.font.size = Pt(11)
            # "인턴 이후 정식 학석사 입학" → 수정
            if '인턴 이후 정식 학석사' in t:
                r.text = '인턴 이후 학석사 연계 가능'

out = 'C:/Users/Ryan/HEART-Lab-Curriculum/안준영_발표자료_수정본.pptx'
prs.save(out)
print('OK:', out)
